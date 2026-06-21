#!/usr/bin/env python3
# coding: utf-8

import os
import uuid
import json
import io
import re
import time
import threading
import traceback
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime, timedelta

import pytz
import requests
import urllib.parse

from flask import (
    Flask,
    request,
    session,
    jsonify,
    redirect,
    render_template_string,
    send_file,
    after_this_request
)

from flask_session import Session
from gtts import gTTS

from superusuarios import (
    es_superusuario,
    obtener_superusuario,
    rol_superusuario,
    nivel_superusuario
)

from usuarios import registrar_usuario, autenticar_usuario
from suscripciones import usuario_premium, aviso_vencimiento
from suscripciones import activar_premium
from profesor_ingles import (
    obtener_perfil as ingles_perfil,
    guardar_perfil as ingles_guardar,
    actualizar_racha as ingles_racha,
    sumar_puntos as ingles_puntos,
    registrar_error as ingles_error,
    completar_leccion as ingles_completar,
    resumen_progreso as ingles_resumen,
    prompt_conversacion_basica,
    prompt_leccion,
    prompt_escenario,
    NIVELES, ESCENARIOS,
)

from openai import OpenAI
from io import BytesIO
from PIL import Image
from docx.shared import Inches
import base64

client = OpenAI()

# --- librerías adicionales para documentos ---
import PyPDF2
from docx import Document as DocxDocument  # para crear / leer .docx
import docx as docx_reader  # para leer .docx (Document ya importado para crear)

# --- librerías adicionales para presentaciones (PPTX) ---
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------------- CONFIG ----------------
APP_NAME = "FOSCHI IA WEB"
CREADOR = "Gustavo Enrique Foschi"
DATA_DIR = "data"
STATIC_DIR = "static"
TEMP_DIR = os.path.join(DATA_DIR, "temp_docs")
IMAGES_DIR = os.path.join(DATA_DIR, "temp_images")

os.makedirs(IMAGES_DIR, exist_ok=True)
os.makedirs(DATA_DIR, exist_ok=True)
os.makedirs(STATIC_DIR, exist_ok=True)
os.makedirs(TEMP_DIR, exist_ok=True)

# ---------------- KEYS ---------------- (usa variables de entorno)
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GOOGLE_CSE_ID = os.getenv("GOOGLE_CSE_ID")
OWM_API_KEY = os.getenv("OWM_API_KEY")

# ---------------- APP ----------------
app = Flask(__name__)
app.secret_key = "FoschiWebKey"
app.config["SESSION_TYPE"] = "filesystem"
Session(app)

# ---------------- UTIL / CACHE / HTTP ----------------
HTTPS = requests.Session()
URL_REGEX = re.compile(r'(https?://[^\s]+)', re.UNICODE)

MEMORY_FILE = os.path.join(DATA_DIR, "memory.json")
MEMORY_CACHE = {}

from datetime import date

def puede_preguntar(usuario):
    if es_superusuario(usuario):
        return True  # 👑 sin límites

    hoy = date.today().isoformat()

    if usuario.get("fecha_preguntas") != hoy:
        usuario["fecha_preguntas"] = hoy
        usuario["preguntas_hoy"] = 0

    if usuario["preguntas_hoy"] >= 5:
        return False

    usuario["preguntas_hoy"] += 1
    return True

def load_json(path):
    """Carga memory.json en cache en RAM para accesos rápidos."""
    global MEMORY_CACHE
    if MEMORY_CACHE:
        return MEMORY_CACHE
    if not os.path.exists(path):
        MEMORY_CACHE = {}
        return MEMORY_CACHE
    try:
        with open(path, "r", encoding="utf-8") as f:
            MEMORY_CACHE = json.load(f)
            return MEMORY_CACHE
    except:
        MEMORY_CACHE = {}
        return MEMORY_CACHE

def save_json(path, data):
    """Guarda MEMORY_CACHE actualizado en disco."""
    global MEMORY_CACHE
    MEMORY_CACHE.update(data)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(MEMORY_CACHE, f, ensure_ascii=False, indent=2)

def fecha_hora_en_es():
    tz = pytz.timezone("America/Argentina/Buenos_Aires")
    ahora = datetime.now(tz)
    meses = ["enero","febrero","marzo","abril","mayo","junio","julio","agosto","septiembre","octubre","noviembre","diciembre"]
    dias = ["lunes","martes","miércoles","jueves","viernes","sábado","domingo"]
    return f"{dias[ahora.weekday()]}, {ahora.day} de {meses[ahora.month-1]} de {ahora.year}, {ahora.hour:02d}:{ahora.minute:02d}"

def hacer_links_clicleables(texto):
    return URL_REGEX.sub(r'<a href="\1" target="_blank" style="color:#ff0000;">\1</a>', texto)

# ---------------- HISTORIAL POR USUARIO ----------------
def guardar_en_historial(usuario, entrada, respuesta):
    path = os.path.join(DATA_DIR, f"{usuario}.json")
    datos = []
    if os.path.exists(path):
        try:
            with open(path,"r",encoding="utf-8") as f:
                datos = json.load(f)
        except:
            datos = []
    datos.append({
        "fecha": datetime.now(pytz.timezone("America/Argentina/Buenos_Aires")).strftime("%d/%m/%Y %H:%M:%S"),
        "usuario": entrada,
        "foschi": respuesta
    })
    datos = datos[-200:]
    try:
        with open(path,"w",encoding="utf-8") as f:
            json.dump(datos,f,ensure_ascii=False,indent=2)
    except Exception as e:
        print("Error guardando historial:", e)

def cargar_historial(usuario):
    path = os.path.join(DATA_DIR, f"{usuario}.json")
    if not os.path.exists(path): return []
    try:
        with open(path,"r",encoding="utf-8") as f:
            return json.load(f)
    except:
        return []

# ---------------- CLIMA ----------------
def obtener_clima(ciudad=None, lat=None, lon=None):
    if not OWM_API_KEY:
        return "No está configurada la API de clima (OWM_API_KEY)."
    try:
        if lat and lon:
            url = f"http://api.openweathermap.org/data/2.5/weather?lat={lat}&lon={lon}&appid={OWM_API_KEY}&units=metric&lang=es"
        else:
            ciudad = ciudad if ciudad else "Buenos Aires"
            url = f"http://api.openweathermap.org/data/2.5/weather?q={ciudad}&appid={OWM_API_KEY}&units=metric&lang=es"
        r = HTTPS.get(url, timeout=3)
        data = r.json()
        if r.status_code != 200:
            msg = data.get("message", "Respuesta no OK de OpenWeatherMap.")
            return f"No pude obtener el clima: {r.status_code} - {msg}"
        desc = data.get("weather", [{}])[0].get("description", "Sin descripción").capitalize()
        temp = data.get("main", {}).get("temp")
        hum = data.get("main", {}).get("humidity")
        name = data.get("name", ciudad if ciudad else "la ubicación")
        parts = [f"El clima en {name} es {desc}"]
        if temp is not None:
            parts.append(f"temperatura {round(temp)}°C")
        if hum is not None:
            parts.append(f"humedad {hum}%")
        return ", ".join(parts) + "."
    except:
        return "No pude obtener el clima."

# ---------------- RECORDATORIOS ----------------
RECORD_FILE = os.path.join(DATA_DIR, "recordatorios.json")
TZ = pytz.timezone("America/Argentina/Buenos_Aires")

def load_recordatorios():
    if not os.path.exists(RECORD_FILE): return []
    try:
        with open(RECORD_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except:
        return []

def save_recordatorios(lista):
    try:
        with open(RECORD_FILE, "w", encoding="utf-8") as f:
            json.dump(lista, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print("Error guardando recordatorios:", e)

def interpretar_fecha_hora(texto):
    """Intenta interpretar frases de tiempo en español. Devuelve datetime (con TZ) o None."""
    ahora = datetime.now(TZ)

    m = re.search(r"en (\d+)\s*minutos?", texto)
    if m:
        return ahora + timedelta(minutes=int(m.group(1)))

    m = re.search(r"en (\d+)\s*horas?", texto)
    if m:
        return ahora + timedelta(hours=int(m.group(1)))

    m = re.search(r"mañana a las (\d{1,2})(?::(\d{2}))?", texto)
    if m:
        hora = int(m.group(1))
        minuto = int(m.group(2)) if m.group(2) else 0
        mañana = (ahora + timedelta(days=1)).replace(hour=hora, minute=minuto, second=0, microsecond=0)
        return mañana

    m = re.search(r"a las (\d{1,2}):(\d{2})", texto)
    if m:
        hora = int(m.group(1))
        minuto = int(m.group(2))
        posible = ahora.replace(hour=hora, minute=minuto, second=0, microsecond=0)
        if posible <= ahora:
            posible = posible + timedelta(days=1)
        return posible

    m = re.search(r"el (\d{1,2}) de (\w+) a las (\d{1,2})(?::(\d{2}))?", texto)
    if m:
        dia = int(m.group(1))
        mes_texto = m.group(2).lower()
        hora = int(m.group(3))
        minuto = int(m.group(4)) if m.group(4) else 0
        meses = {
            "enero":1,"febrero":2,"marzo":3,"abril":4,"mayo":5,"junio":6,
            "julio":7,"agosto":8,"septiembre":9,"octubre":10,"noviembre":11,"diciembre":12
        }
        mes = meses.get(mes_texto)
        if mes:
            año = ahora.year
            try:
                candidato = datetime(año, mes, dia, hora, minuto)
                candidato = TZ.localize(candidato)
            except Exception:
                return None
            if candidato <= ahora:
                try:
                    candidato = datetime(año+1, mes, dia, hora, minuto)
                    candidato = TZ.localize(candidato)
                except:
                    return None
            return candidato

    return None

def agregar_recordatorio(usuario, motivo_texto, fecha_hora_dt):
    """Agrega un recordatorio persistente. fecha_hora_dt debe ser datetime con TZ (o naive en TZ)."""
    if fecha_hora_dt.tzinfo is None:
        fecha_hora_dt = fecha_hora_dt.replace(tzinfo=TZ)
    lista = load_recordatorios()
    lista.append({
        "usuario": usuario,
        "motivo": motivo_texto.strip(),
        "cuando": fecha_hora_dt.strftime("%Y-%m-%d %H:%M:%S")
    })
    save_recordatorios(lista)

def listar_recordatorios(usuario):
    lista = load_recordatorios()
    return [r for r in lista if r.get("usuario") == usuario]

def borrar_recordatorios(usuario):
    lista = load_recordatorios()
    lista = [r for r in lista if r.get("usuario") != usuario]
    save_recordatorios(lista)

def monitor_recordatorios():
    """Hilo que revisa recordatorios y los dispara. Cuando se dispara, lo guarda en el historial del usuario."""
    while True:
        try:
            lista = load_recordatorios()
            ahora = datetime.now(TZ)
            restantes = []
            for r in lista:
                try:
                    cuando = datetime.strptime(r["cuando"], "%Y-%m-%d %H:%M:%S")
                    cuando = TZ.localize(cuando)
                except Exception:
                    try:
                        cuando = datetime.fromisoformat(r["cuando"])
                        if cuando.tzinfo is None:
                            cuando = TZ.localize(cuando)
                    except:
                        continue
                if cuando <= ahora:
                    usuario = r.get("usuario", "anon")
                    motivo = r.get("motivo", "(sin motivo)")
                    aviso_texto = f"⏰ Tenés un recordatorio: {motivo}"
                    try:
                        guardar_en_historial(usuario, f"[recordatorio] {motivo}", aviso_texto)
                    except Exception:
                        pass
                    print(aviso_texto)
                else:
                    restantes.append(r)
            save_recordatorios(restantes)
        except Exception as e:
            print("Error en monitor_recordatorios:", e)
        time.sleep(30)

# iniciar hilo del monitor (daemon)
threading.Thread(target=monitor_recordatorios, daemon=True).start()

# ---------------- learn_from_message (registro de memoria) ----------------
def learn_from_message(usuario, mensaje, respuesta):
    try:
        memory = load_json(MEMORY_FILE)
        if usuario not in memory:
            memory[usuario] = {"temas": {}, "mensajes": [], "ultima_interaccion": None}
        # Guardar texto en memoria (limitamos)
        memory[usuario]["mensajes"].append({"usuario": str(mensaje), "foschi": str(respuesta)})
        memory[usuario]["mensajes"] = memory[usuario]["mensajes"][-200:]
        ahora = datetime.now(pytz.timezone("America/Argentina/Buenos_Aires"))
        memory[usuario]["ultima_interaccion"] = ahora.strftime("%d/%m/%Y %H:%M:%S")
        # Tópicos simples
        for palabra in str(mensaje).lower().split():
            if len(palabra) > 3:
                memory[usuario]["temas"][palabra] = memory[usuario]["temas"].get(palabra, 0) + 1
        save_json(MEMORY_FILE, memory)
    except Exception as e:
        print("Error en learn_from_message:", e)

# ---------------- DICTADO A WORD ----------------
@app.route("/dictado_word", methods=["POST"])
def dictado_word():

    data = request.get_json(silent=True) or {}

    texto = data.get("texto", "").strip()

    if not texto:
        return jsonify({
            "ok": False,
            "error": "Texto vacío"
        })

    nombre = f"dictado_{uuid.uuid4().hex}.docx"

    ruta = os.path.join(TEMP_DIR, nombre)

    doc = DocxDocument()

    doc.add_heading("Dictado Foschi IA", 0)

    for linea in texto.split("\n"):

        linea = linea.strip()

        if linea:
            doc.add_paragraph(linea)

    doc.save(ruta)

    @after_this_request
    def remove_file(response):

        try:
            if os.path.exists(ruta):
                os.remove(ruta)
        except Exception as e:
            print("Error eliminando temporal:", e)

        return response

    return send_file(
        ruta,
        as_attachment=True,
        download_name="dictado_foschi.docx"
    )
@app.route(
    "/imagen_a_word",
    methods=["POST"]
)
def imagen_a_word():

    if "imagen" not in request.files:
        return "No se recibió imagen",400

    archivo = request.files["imagen"]

    nombre = (
        uuid.uuid4().hex +
        ".png"
    )

    ruta_imagen = os.path.join(
        IMAGES_DIR,
        nombre
    )

    archivo.save(ruta_imagen)

    try:

        with open(
            ruta_imagen,
            "rb"
        ) as f:

            imagen_base64 = base64.b64encode(
                f.read()
            ).decode()

        respuesta = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role":"user",
                    "content":[
                        {
                            "type":"text",
                            "text":"Extraé TODO el texto visible."
                        },
                        {
                            "type":"image_url",
                            "image_url":{
                                "url":"data:image/png;base64," + imagen_base64
                            }
                        }
                    ]
                }
            ],
            max_tokens=4000
        )

        texto = respuesta.choices[0].message.content

        doc = DocxDocument()

        doc.add_heading(
            "Documento extraído",
            0
        )

        doc.add_picture(
            ruta_imagen,
            width=Inches(4)
        )

        doc.add_paragraph(texto)

        salida = os.path.join(
            TEMP_DIR,
            "ocr_" +
            uuid.uuid4().hex +
            ".docx"
        )

        doc.save(salida)

        @after_this_request
        def remove_file(response):

            try:

                if os.path.exists(salida):
                    os.remove(salida)

                if os.path.exists(ruta_imagen):
                    os.remove(ruta_imagen)

            except Exception as e:
                print("Error eliminando temporales:", e)

            return response

        return send_file(
            salida,
            as_attachment=True,
            download_name="imagen_extraida.docx"
        )

    except Exception as e:

        return str(e),500
      
@app.route(
    "/editar_imagen",
    methods=["POST"]
)
def editar_imagen():

    try:

        if "imagen" not in request.files:
            return jsonify({
                "ok": False,
                "error": "No se recibió imagen"
            }), 400

        imagen = request.files["imagen"]

        print("NOMBRE:", imagen.filename)
        print("TIPO:", imagen.content_type)

        contenido = BytesIO(
            imagen.read()
        )

        contenido.name = imagen.filename

        # "quality" alto da más detalle pero tarda mucho más y puede
        # provocar timeouts del servidor/proxy. "medium" es un buen
        # equilibrio; podés probar "high" si tu hosting lo soporta.
        resultado = client.images.edit(
            model="gpt-image-1",
            image=contenido,
            prompt=request.form.get(
                "prompt",
                ""
            ),
            size="1024x1024",
            quality="medium"
        )

        return jsonify({
            "ok": True,
            "imagen": resultado.data[0].b64_json
        })

    except Exception as e:

        print("ERROR EDITAR IMAGEN:")
        print(str(e))

        return jsonify({
            "ok": False,
            "error": str(e)
        }), 500


@app.route(
    "/generar_imagen",
    methods=["POST"]
)
def generar_imagen():

    try:

        data = request.get_json(silent=True) or {}
        prompt = (data.get("prompt") or "").strip()

        if not prompt:
            return jsonify({
                "ok": False,
                "error": "No se recibió descripción para generar la imagen"
            }), 400

        # "quality" alto da más detalle pero tarda mucho más y puede
        # provocar timeouts del servidor/proxy. "medium" es un buen
        # equilibrio; podés probar "high" si tu hosting lo soporta.
        resultado = client.images.generate(
            model="gpt-image-1",
            prompt=prompt,
            size="1024x1024",
            quality="medium"
        )

        return jsonify({
            "ok": True,
            "imagen": resultado.data[0].b64_json
        })

    except Exception as e:

        print("ERROR GENERAR IMAGEN:")
        print(str(e))

        return jsonify({
            "ok": False,
            "error": str(e)
        }), 500

# ---------------- RESPUESTA IA ----------------

def generar_respuesta(mensaje, usuario, lat=None, lon=None, tz=None, max_hist=5):
       
    # Bloqueo por no premium
    if not usuario_premium(usuario) and not es_superusuario(usuario):
        if len(mensaje) > 200:
            return {
                "texto": "🔒 Esta función es solo para usuarios Premium.\n\n💎 Activá Foschi IA Premium desde el botón superior para seguir.",
                "imagenes": [],
                "borrar_historial": False
            }

    # Asegurar string
    if not isinstance(mensaje, str):
        mensaje = str(mensaje)

    mensaje_lower = mensaje.lower().strip()
           
    # --- RECORDATORIOS: comandos y detección ---
    try:
        if mensaje_lower in ["mis recordatorios", "lista de recordatorios", "ver recordatorios"]:
            recs = listar_recordatorios(usuario)
            if not recs:
                return {"texto": "📭 No tenés recordatorios pendientes.", "imagenes": [], "borrar_historial": False}
            texto = "📌 Tus recordatorios:\n" + "\n".join([f"- {r['motivo']} → {r['cuando']}" for r in recs])
            return {"texto": texto, "imagenes": [], "borrar_historial": False}

        if "borrar recordatorios" in mensaje_lower or "eliminar recordatorios" in mensaje_lower:
            borrar_recordatorios(usuario)
            return {"texto": "🗑️ Listo, eliminé todos tus recordatorios.", "imagenes": [], "borrar_historial": False}

        if mensaje_lower.startswith(("recordame", "haceme acordar", "avisame", "recordá")):
            fecha_hora = interpretar_fecha_hora(mensaje_lower)
            if fecha_hora is None:
                return {"texto": "⏰ Decime cuándo: ejemplo 'mañana a las 9', 'en 15 minutos' o 'el 5 de diciembre a las 18'.", "imagenes": [], "borrar_historial": False}
            motivo = mensaje
            for p in ["recordame", "haceme acordar", "avisame", "recordá"]:
                motivo = re.sub(p, "", motivo, flags=re.IGNORECASE).strip()
            if not motivo:
                motivo = "Recordatorio"
            agregar_recordatorio(usuario, motivo, fecha_hora)
            return {"texto": f"✅ Listo, te lo recuerdo el {fecha_hora.strftime('%d/%m %H:%M')}.", "imagenes": [], "borrar_historial": False}
    except Exception as e:
        print("Error en manejo de recordatorios:", e)

    # BORRAR HISTORIAL
    if any(p in mensaje_lower for p in ["borrar historial", "limpiar historial", "reset historial"]):
        path = os.path.join(DATA_DIR, f"{usuario}.json")
        if os.path.exists(path): os.remove(path)
        memory = load_json(MEMORY_FILE)
        if usuario in memory:
            memory[usuario]["mensajes"] = []
            save_json(MEMORY_FILE, memory)
        return {"texto": "✅ Historial borrado correctamente.", "imagenes": [], "borrar_historial": True}

    # FECHA / HORA
    if any(p in mensaje_lower for p in ["qué día", "que día", "qué fecha", "que fecha", "qué hora", "que hora", "día es hoy", "fecha hoy"]):
        texto = fecha_hora_en_es()
        learn_from_message(usuario, mensaje, texto)
        return {"texto": texto, "imagenes": [], "borrar_historial": False}

    # CLIMA
    if "clima" in mensaje_lower:
        ciudad_match = re.search(r"clima en ([a-zA-ZáéíóúÁÉÍÓÚñÑ\s]+)", mensaje_lower)
        ciudad = ciudad_match.group(1).strip() if ciudad_match else None
        texto = obtener_clima(ciudad=ciudad, lat=lat, lon=lon)
        learn_from_message(usuario, mensaje, texto)
        return {"texto": texto, "imagenes": [], "borrar_historial": False}

    # INFORMACIÓN ACTUALIZADA (NOTICIAS)
    if any(word in mensaje_lower for word in ["presidente", "actualidad", "noticias", "quién es", "últimas noticias", "evento actual"]):
        resultados = []
        if GOOGLE_API_KEY and GOOGLE_CSE_ID:
            try:
                url = (
                    f"https://www.googleapis.com/customsearch/v1"
                    f"?key={GOOGLE_API_KEY}&cx={GOOGLE_CSE_ID}"
                    f"&q={urllib.parse.quote(mensaje)}&sort=date"
                )
                r = HTTPS.get(url, timeout=5)
                data = r.json()
                for item in data.get("items", [])[:5]:
                    snippet = item.get("snippet", "").strip()
                    if snippet and snippet not in resultados:
                        resultados.append(snippet)
            except Exception as e:
                print("Error al obtener noticias:", e)

        if resultados:
            texto_bruto = " ".join(resultados)
            try:
                client = OpenAI(api_key=OPENAI_API_KEY)
                prompt = (
                    f"Tengo estos fragmentos de texto recientes: {texto_bruto}\n\n"
                    f"Respondé a la pregunta: '{mensaje}'. "
                    f"Usá un tono natural y directo en español argentino, sin frases como "
                    f"'según los textos', 'según los fragmentos' o 'de acuerdo a las fuentes'. "
                    f"Contestá con una sola oración clara y actualizada. Si no hay información suficiente, decílo sin inventar."
                )

                resp = client.chat.completions.create(
                    model="gpt-4-turbo",
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.5,
                    max_tokens=120
                )

                texto = resp.choices[0].message.content.strip()
            except Exception as e:
                texto = "No pude generar la respuesta con OpenAI."
        else:
            texto = "No pude obtener información actualizada en este momento."

        learn_from_message(usuario, mensaje, texto)
        return {"texto": texto, "imagenes": [], "borrar_historial": False}

    # QUIÉN CREÓ / PREGUNTAS ESTÁTICAS
    if any(p in mensaje_lower for p in [
        "quién te creó", "quien te creo",
        "quién te hizo", "quien te hizo",
        "quién te programó", "quien te programo",
        "quién te inventó", "quien te invento",
        "quién te desarrolló", "quien te desarrollo",
        "quién te construyó", "quien te construyo"
    ]):
        texto = "Fui creada por Gustavo Enrique Foschi, el mejor 😎."
        learn_from_message(usuario, mensaje, texto)
        return {"texto": texto, "imagenes": [], "borrar_historial": False}

    # RESULTADOS DEPORTIVOS (actualizados)
    if any(p in mensaje_lower for p in [
        "resultado", "marcador", "ganó", "empató", "perdió",
        "partido", "deporte", "fútbol", "futbol", "nba", "tenis", "f1", "formula 1", "motogp"
    ]):
        resultados = []
        if GOOGLE_API_KEY and GOOGLE_CSE_ID:
            try:
                url = (
                    f"https://www.googleapis.com/customsearch/v1"
                    f"?key={GOOGLE_API_KEY}&cx={GOOGLE_CSE_ID}"
                    f"&q={urllib.parse.quote(mensaje + ' resultados deportivos actualizados')}"
                    f"&sort=date"
                )
                r = HTTPS.get(url, timeout=5)
                data = r.json()
                for item in data.get("items", [])[:5]:
                    snippet = item.get("snippet", "").strip()
                    if snippet and snippet not in resultados:
                        resultados.append(snippet)
            except Exception as e:
                print("Error al obtener resultados deportivos:", e)

        if resultados:
            texto_bruto = " ".join(resultados)
            try:
                client = OpenAI(api_key=OPENAI_API_KEY)
                prompt = (
                    f"Tengo estos fragmentos recientes sobre deportes: {texto_bruto}\n\n"
                    f"Respondé brevemente la consulta '{mensaje}' con los resultados deportivos actuales. "
                    f"Usá un tono natural, tipo boletín deportivo argentino, sin frases como 'según los textos'. "
                    f"Respondé en una sola oración clara."
                )

                resp = client.chat.completions.create(
                    model="gpt-4-turbo",
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.5,
                    max_tokens=150
                )
                texto = resp.choices[0].message.content.strip()
            except Exception as e:
                texto = "No pude generar la respuesta deportiva."
        else:
            texto = "No pude encontrar resultados deportivos recientes en este momento."

        learn_from_message(usuario, mensaje, texto)
        return {"texto": texto, "imagenes": [], "borrar_historial": False}

    # PRESENTACIONES (POWERPOINT)
    if any(p in mensaje_lower for p in [
        "presentación", "presentacion", "powerpoint", "power point",
        "diapositivas", "diapositiva", "slides", "ppt"
    ]):
        tema_pre = (
            mensaje
            .replace("\\", "\\\\")
            .replace("'", "\\'")
            .replace('"', "&quot;")
            .replace("\n", " ")
            .strip()
        )
        texto = (
            "🖥️ ¡Puedo armarte una presentación de PowerPoint (.pptx) desde cero, "
            "con imágenes generadas por IA o videos cortos que vos subas! "
            "También puedo basarme en un documento que ya hayas subido. "
            "Tocá el botón para configurarla:<br><br>"
            f"<button onclick=\"abrirGeneradorPresentacion('{tema_pre}')\">🖥️ Crear presentación</button>"
        )
        learn_from_message(usuario, mensaje, texto)
        return {"texto": texto, "imagenes": [], "borrar_historial": False}

    # SALIDA GENERAL: pasar a OpenAI para respuesta conversacional
    try:
        memoria = load_json(MEMORY_FILE)
        historial = memoria.get(usuario, {}).get("mensajes", [])[-max_hist:]
        resumen = " ".join([m["usuario"] + ": " + m["foschi"] for m in historial[-3:]]) if historial else ""

        client = OpenAI(api_key=OPENAI_API_KEY)

        prompt_messages = [
            {
                "role": "system",
                "content": (
                    "Sos FOSCHI IA, una inteligencia amable, directa y con humor ligero. "
                    "Tus respuestas deben ser claras, ordenadas y sonar naturales en español argentino. "
                    "Si el usuario pide información o ayuda técnica, explicá paso a paso y sin mezclar temas. "
                    f"Resumen de últimas interacciones: {resumen if resumen else 'ninguna.'}"
                )
            },
            {"role": "user", "content": mensaje}
        ]

        resp = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=prompt_messages,
            temperature=0.7,
            max_tokens=700
        )

        texto = resp.choices[0].message.content.strip()

    except Exception as e:
        texto = f"No pude generar respuesta: {e}"

        texto = hacer_links_clicleables(texto)

    aviso = aviso_vencimiento(usuario)
    if aviso:
        texto += "\n\n" + aviso

    learn_from_message(usuario, mensaje, texto)
    return {
        "texto": texto,
        "imagenes": [],
        "borrar_historial": False
    }

# ---------------- Plantilla HTML (modificada para menu clip + subir pdf/docx) ----------------
HTML_TEMPLATE = """  
<!doctype html>
<html>
<head>
<title>{{APP_NAME}}</title>
<meta name="viewport" content="width=device-width, initial-scale=1">
<style>
/* --- ESTILOS GENERALES --- */
body{
 font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
 background:#000814;
 color:#00eaff;
 margin:0;
 padding:0;
 display:flex;
 flex-direction:column;
 height:100vh;
 overflow:hidden;
 text-shadow:0 0 6px #00eaff;
}

/* --- HEADER SUPERIOR (LOGO + BOTONES) --- */
#header{
  display:flex;
  align-items:center;
  justify-content:space-between;
  padding:8px 16px;
  background: linear-gradient(#000814,#00111a);
  flex-shrink:0;
  position:sticky;
  top:0;
  z-index:10;
  box-shadow:0 0 12px #00eaff66;
}

#leftButtons{
  display:flex;
  align-items:center;
  gap:12px;
}

#rightButtons{
  display:flex;
  align-items:center;
  gap:8px;
}

.separator{
  border-left:1px solid #00eaff55;
  height:32px;
  margin:0 8px;
}

#header button{
  font-size:14px;
  padding:6px 10px;
  border-radius:6px;
  background:#001f2e;
  color:#00eaff;
  border:1px solid #006688;
  cursor:pointer;
  text-shadow:0 0 4px #00eaff;
  box-shadow:0 0 8px #0099bb;
  transition:0.3s;
}
/* ===================== */
/* === MODO DAY ======== */
/* ===================== */
body.day{
  background:#ffffff;
  color:#000000;
  text-shadow:none;
}

body.day #header{
  background:#ffffff;
  box-shadow:0 2px 8px rgba(0,0,0,0.1);
}

body.day #chat{
  background:#f5f5f5;
  border-top:1px solid #ccc;
  border-bottom:1px solid #ccc;
  box-shadow:none;
}

body.day #inputBar{
  background:#ffffff;
  border-top:1px solid #ccc;
}

body.day button{
  background:#ffffff !important;
  color:#000000 !important;
  border:1px solid #000000 !important;
  box-shadow:none !important;
  text-shadow:none !important;
}

body.day input{
  background:#ffffff !important;
  color:#000000 !important;
  border:1px solid #000000 !important;
  box-shadow:none !important;
}

/* Premium en Day */
body.day #premiumBtn{
  animation:none;
  font-weight:600;
}

#header button:hover{
  background:#003547;
  box-shadow:0 0 14px #00eaff;
}

#premiumBtn{
  animation: neonGlow 1.5s ease-in-out infinite alternate;
}

@keyframes neonGlow {
  0% { box-shadow: 0 0 8px #00eaff, 0 0 12px #00eaff, 0 0 16px #00eaff; color:#00eaff; }
  50% { box-shadow: 0 0 12px #00ffff, 0 0 20px #00ffff, 0 0 28px #00ffff; color:#00ffff; }
  100% { box-shadow: 0 0 8px #00eaff, 0 0 12px #00eaff, 0 0 16px #00eaff; color:#00eaff; }
}

#logo{
  width:120px;
  cursor:pointer;
  transition: transform 0.5s, filter 0.5s;
  filter: drop-shadow(0 0 12px #00eaff);
}
#logo:hover{
  transform:scale(1.2) rotate(6deg);
  filter:drop-shadow(0 0 20px #00eaff);
}

/* --- CHAT --- */
#chat{
  flex:1;
  overflow-y:auto;
  padding:10px;
  background: linear-gradient(#00111a,#000814);
  border-top:2px solid #00eaff44;
  border-bottom:2px solid #00eaff44;
  box-shadow: inset 0 0 15px #00eaff55;
  padding-bottom:120px; /* espacio para la barra inferior */
}

/* --- MENSAJES --- */
.message{
 margin:5px 0;
 padding:8px 12px;
 border-radius:15px;
 max-width:80%;
 word-wrap:break-word;
 opacity:0;
 transition:opacity 0.5s, box-shadow 0.5s, background 0.5s;
 font-size:15px;
}
.message.show{ opacity:1; }
.user{
 background:rgba(51,0,255,0.3);
 color:#b4b7ff;
 margin-left:auto;
 text-align:right;
 border:1px solid #4455ff;
 box-shadow:0 0 8px #3344ff;
}
.ai{
 background:rgba(0,255,255,0.2);
 color:#00eaff;
 margin-right:auto;
 text-align:left;
 border:1px solid #00eaff;
 box-shadow:0 0 10px #00eaff;
}
a{ color:#00eaff; text-decoration:underline; }
img{ max-width:300px; border-radius:10px; margin:5px 0; box-shadow:0 0 10px #00eaff88; border:1px solid #00eaff55; }

/* --- BARRA DE ENTRADA FIJA ABAJO --- */
#inputBar{
 display:flex;
 align-items:center;
 gap:6px;
 padding:8px;
 background:#001d29;
 border-top:2px solid #00eaff44;
 flex-shrink:0;
 position:fixed;
 bottom:0;
 left:0;
 width:100%;
 box-sizing:border-box;
 z-index:20;
}
#inputBar input[type=text]{
 flex:1;
 padding:10px;
 font-size:16px;
 background:#00121d;
 color:#00eaff;
 border:1px solid #003344;
 box-shadow:0 0 6px #00eaff55 inset;
 border-radius:5px;
}
#inputBar button{
 padding:10px;
 font-size:16px;
 border:none;
 border-radius:5px;
 background:#001f2e;
 color:#00eaff;
 cursor:pointer;
 border:1px solid #006688;
 text-shadow:0 0 4px #00eaff;
 box-shadow:0 0 8px #0099bb;
 transition:0.25s;
}
#inputBar button:hover{
 background:#003547;
 box-shadow:0 0 14px #00eaff;
}

/* --- BOTONES PEQUEÑOS --- */
#vozBtn,#borrarBtn,#premiumBtn{ font-size:14px; padding:6px 10px; }

/* --- BOTÓN ADJUNTAR (CLIP) --- */
#clipBtn{
  width:40px;
  height:40px;
  display:flex;
  align-items:center;
  justify-content:center;
  font-size:20px;
  cursor:pointer;
  border-radius:6px;
  background:#001f2e;
  color:#00eaff;
  border:1px solid #006688;
  box-shadow:0 0 10px #0099bb;
  transition:0.25s;
  user-select:none;
}

#clipBtn:hover{
  background:#003547;
  box-shadow:0 0 16px #00eaff;
  transform:scale(1.08);
}

/* --- CLIP EN MODO DAY --- */
body.day #clipBtn{
  background:#ffffff;
  color:#000000;
  border:1px solid #000000;
  box-shadow:none;
}

body.day #clipBtn:hover{
  background:#f0f0f0;
}

/* --- MENÚ DE ADJUNTOS (overlay fijo sobre la barra) --- */
#adjuntos_menu{
 position:fixed;
 left:10px;
 bottom:62px;
 display:none;
 background:#001f2e;
 border:1px solid #003547;
 padding:8px;
 border-radius:8px;
 box-shadow:0 6px 16px rgba(0,0,0,0.6);
 z-index:999;
 max-width:90vw;
}
#adjuntos_menu button{ display:block; width:200px; margin:4px 0; text-align:left; transition:0.2s; }
.hidden_file_input{ display:none; }

/* --- MENÚ PREMIUM --- */
#premiumMenu button{
  display:block;
  width:120px;
  margin:4px 0;
  text-align:left;
}

/* Opciones del menú de adjuntos: se iluminan al pasar el mouse o al elegirlas */
#adjuntos_menu button:hover,
#adjuntos_menu button.opcion-seleccionada{
  background:#003547 !important;
  color:#ffd700 !important;
  border:1px solid #ffd700 !important;
  box-shadow:0 0 16px #ffd700, 0 0 28px #ffd70066 !important;
}

#adjuntos_menu button.opcion-seleccionada{
  animation: iluminarOpcionAdjunto 0.4s ease-in-out;
}

@keyframes iluminarOpcionAdjunto{
  0%{ box-shadow:0 0 2px #ffd70033; }
  50%{ box-shadow:0 0 26px #ffd700; }
  100%{ box-shadow:0 0 16px #ffd700; }
}

/* --- AVISO "ACÁ ESTÁ EL MENÚ" SOBRE EL CLIP (solo usuarios no premium) --- */
#clipBtn{ position:relative; }

#adjuntosHint{
  position:absolute;
  bottom:50px;
  left:0;
  background:#001f2e;
  color:#00eaff;
  border:1px solid #00eaff;
  font-size:12px;
  padding:6px 10px;
  border-radius:6px;
  white-space:nowrap;
  box-shadow:0 0 14px #00eaff99;
  z-index:99;
  pointer-events:none;
  animation: pulsoAdjuntosHint 1.6s ease-in-out infinite;
}

#adjuntosHint::after{
  content:"";
  position:absolute;
  bottom:-6px;
  left:14px;
  border-width:6px 6px 0 6px;
  border-style:solid;
  border-color:#00eaff transparent transparent transparent;
}

@keyframes pulsoAdjuntosHint{
  0%,100%{ box-shadow:0 0 8px #00eaff66; opacity:0.9; }
  50%{ box-shadow:0 0 20px #00eaff; opacity:1; }
}

body.day #adjuntosHint{
  background:#ffffff;
  color:#0077aa;
  border-color:#0077aa;
}
body.day #adjuntosHint::after{
  border-color:#0077aa transparent transparent transparent;
}

/* --- AJUSTES RESPONSIVE PARA MÓVIL --- */
@media (max-width:600px){
  #inputBar input[type=text]{ font-size:18px; padding:12px; }
  #inputBar button{ font-size:16px; padding:10px; }
  #logo{ width:140px; } /* logo más grande en móvil */
}
/* =============================== */
/* === MENSAJES EN MODO DAY ====== */
/* =============================== */

body.day .user{
  background:#e9e9e9;
  color:#000000;
  border:1px solid #000000;
  box-shadow:none;
}

body.day .ai{
  background:#f5f5f5;
  color:#000000;
  border:1px solid #000000;
  box-shadow:none;
}

body.day .ai a,
body.day .user a{
  color:#000000;
}

/* ===== FORMULARIO PREMIUM CENTRADO ===== */
.form-premium{
  width:100%;
  max-width:320px;
  margin:0 auto;
  display:flex;
  flex-direction:column;
  align-items:center;
  gap:10px;
}

.form-premium input{
  width:100%;
  padding:10px;
  border-radius:6px;
  border:1px solid #006688;
  background:#00121d;
  color:#00eaff;
  box-shadow:0 0 6px #00eaff55 inset;
}

.form-premium button{
  width:100%;
  padding:10px;
  border-radius:6px;
  border:1px solid #006688;
  background:#001f2e;
  color:#00eaff;
  cursor:pointer;
  box-shadow:0 0 8px #0099bb;
}

.form-premium button:hover{
  background:#003547;
  box-shadow:0 0 14px #00eaff;
}

.wave{
 width:6px;
 height:20px;
 background:#00eaff;
 animation:wave 1s infinite ease-in-out;
 border-radius:4px;
}

.wave:nth-child(2){animation-delay:0.1s;}
.wave:nth-child(3){animation-delay:0.2s;}
.wave:nth-child(4){animation-delay:0.3s;}
.wave:nth-child(5){animation-delay:0.4s;}
/* ========================= */
/* 🎤 DICTADO PRO */
/* ========================= */

/* Panel de texto acumulado durante dictado */
#dictadoPanel{
  position:fixed;
  bottom:62px;
  left:0; right:0;
  max-height:180px;
  overflow-y:auto;
  background:#00060d;
  border-top:2px solid #ff003366;
  padding:10px 16px;
  font-size:15px;
  color:#ffffff;
  line-height:1.7;
  z-index:25;
  display:none;
  white-space:pre-wrap;
  word-break:break-word;
  scrollbar-width:thin;
  scrollbar-color:#ff003366 transparent;
}
#dictadoPanel::-webkit-scrollbar{ width:4px; }
#dictadoPanel::-webkit-scrollbar-thumb{ background:#ff003366; border-radius:4px; }

body.day #dictadoPanel{
  background:#f5f5f5;
  color:#111;
  border-color:#cc000055;
}

#dictadoBtn.activo{
  background:#ff0033 !important;
  color:#fff !important;
  border:1px solid #ff3355 !important;
  box-shadow:0 0 20px #ff0033 !important;
  animation:pulseDictado 1s infinite;
}

@keyframes pulseDictado{
  0%{transform:scale(1);}
  50%{transform:scale(1.06);}
  100%{transform:scale(1);}
}

#dictadoEstado{
  box-shadow:0 0 20px #ff0033;
  animation:pulseEstado 1s infinite;
}

@keyframes pulseEstado{
  0%{opacity:1;}
  50%{opacity:0.6;}
  100%{opacity:1;}
}

@keyframes wave{
0%,100%{height:10px;}
50%{height:40px;}
}

</style>
</head>

<body>
<!-- HEADER -->
<div id="header">
  <div id="leftButtons">
    <img src="/static/logo.png" id="logo" onclick="logoClick()" alt="logo">

    <div id="premiumContainer" style="position:relative; margin-left:12px;">
      <button id="dayNightBtn" onclick="toggleDayNight()">🌙</button>

      <button id="premiumBtn" onclick="togglePremiumMenu()">
        {% if premium %}
          💎 Premium activo
        {% else %}
          💎 Pasar a Premium
        {% endif %}
      </button>

      {% if not premium %}
      <div id="premiumMenu"
           style="display:none; position:absolute; top:36px; left:0;
                  background:#001f2e; border:1px solid #003547;
                  border-radius:6px; padding:6px;
                  box-shadow:0 6px 16px rgba(0,0,0,0.6); z-index:100;">
        <button onclick="irPremium('mensual')">💎 Pago Mensual</button>
        <button onclick="irPremium('anual')">💎 Pago Anual</button>
      </div>
      {% endif %}
    </div>
  </div>

  <div id="rightButtons">
    <div class="separator"></div>
    <button onclick="detenerVoz()">⏹️ Detener voz</button>
    <button id="vozBtn" onclick="toggleVoz()">🔊 Voz activada</button>
    <button id="borrarBtn" onclick="borrarPantalla()">🧹 Borrar pantalla</button>
    <button onclick="verHistorial()">🗂️ Historial</button>
  </div>
</div>

<!-- CHAT -->
<div id="chat" role="log" aria-live="polite"></div>

<!-- MODAL FLOTANTE EDITOR DE IMAGEN -->
<div id="editorImagen" style="display:none;position:fixed;inset:0;z-index:9998;background:rgba(0,8,20,0.88);align-items:center;justify-content:center;padding:12px;box-sizing:border-box;">
  <div style="background:linear-gradient(135deg,#001a2e,#002a44);border:1px solid #00eaff55;border-radius:18px;box-shadow:0 0 40px #00eaff33;width:100%;max-width:900px;max-height:94vh;overflow-y:auto;padding:22px;box-sizing:border-box;">

    <!-- Título y botón cerrar -->
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:16px;">
      <span style="color:#00eaff;font-size:17px;font-weight:700;text-shadow:0 0 8px #00eaff;letter-spacing:1px;">✏️ Editor de Imagen IA</span>
      <button onclick="cancelarEdicionImagen()" style="background:transparent;border:1px solid #ff444466;color:#ff6666;border-radius:8px;padding:6px 16px;font-size:15px;cursor:pointer;">✕ Cerrar</button>
    </div>

    <!-- Original | Resultado -->
    <div style="display:flex;gap:12px;flex-wrap:wrap;margin-bottom:14px;">
      <div style="flex:1;min-width:200px;text-align:center;">
        <div style="color:#00eaff66;font-size:11px;letter-spacing:1px;text-transform:uppercase;margin-bottom:6px;">Original</div>
        <img id="previewImagen" src="" class="img-box" style="display:none;">
        <div id="placeholderOrig" class="img-placeholder" style="display:flex;">Sin imagen</div>
      </div>
      <div style="flex:1;min-width:200px;text-align:center;">
        <div style="color:#00eaff66;font-size:11px;letter-spacing:1px;text-transform:uppercase;margin-bottom:6px;">Resultado</div>
        <img id="resultadoImagen" src="" class="img-box" style="display:none;">
        <div id="placeholderResult" class="img-placeholder" style="display:flex;">Aquí aparecerá la edición</div>
      </div>
    </div>

    <!-- Prompt -->
    <textarea id="promptImagen" placeholder="Describí qué querés cambiar... (ej: cambiá el fondo por un bosque, poné lentes al personaje)" style="width:100%;height:75px;background:#001122;color:#00eaff;border:1px solid #006688;border-radius:10px;padding:10px 12px;font-size:14px;resize:vertical;box-sizing:border-box;outline:none;font-family:'Segoe UI',sans-serif;"></textarea>

    <!-- Botones -->
    <div style="display:flex;gap:10px;margin-top:14px;flex-wrap:wrap;align-items:center;">
      <button id="btnAplicarEdicion" type="button" onclick="editarImagenActual()" style="padding:11px 24px;background:linear-gradient(135deg,#005577,#007799);color:#fff;border:none;border-radius:10px;font-size:15px;font-weight:600;cursor:pointer;box-shadow:0 0 14px #00eaff44;">✨ Aplicar con IA</button>
      <a id="btnDescargarEdicion" download="foschi_editada.png" style="display:none;padding:11px 18px;background:linear-gradient(135deg,#004400,#006600);color:#00ff88;border:1px solid #00ff8844;border-radius:10px;font-size:14px;font-weight:600;text-decoration:none;cursor:pointer;">⬇️ Descargar</a>
      <span id="editorEstado" style="display:none;align-items:center;gap:8px;color:#00eaff88;font-size:13px;">
        <span style="display:inline-block;width:15px;height:15px;border:2px solid #00eaff33;border-top-color:#00eaff;border-radius:50%;animation:spinImg 0.7s linear infinite;"></span>
        Procesando con IA...
      </span>
    </div>

  </div>
</div>

<!-- MODAL FLOTANTE GENERADOR DE IMAGEN -->
<div id="generadorImagen" style="display:none;position:fixed;inset:0;z-index:9998;background:rgba(0,8,20,0.88);align-items:center;justify-content:center;padding:12px;box-sizing:border-box;">
  <div style="background:linear-gradient(135deg,#001a2e,#002a44);border:1px solid #00eaff55;border-radius:18px;box-shadow:0 0 40px #00eaff33;width:100%;max-width:900px;max-height:94vh;overflow-y:auto;padding:22px;box-sizing:border-box;">

    <!-- Título y botón cerrar -->
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:16px;">
      <span style="color:#00eaff;font-size:17px;font-weight:700;text-shadow:0 0 8px #00eaff;letter-spacing:1px;">🎨 Generador de Imagen IA</span>
      <button onclick="cancelarGeneradorImagen()" style="background:transparent;border:1px solid #ff444466;color:#ff6666;border-radius:8px;padding:6px 16px;font-size:15px;cursor:pointer;">✕ Cerrar</button>
    </div>

    <!-- Resultado -->
    <div style="text-align:center;margin-bottom:14px;">
      <div style="color:#00eaff66;font-size:11px;letter-spacing:1px;text-transform:uppercase;margin-bottom:6px;">Resultado</div>
      <img id="resultadoGenerador" src="" class="img-box" style="display:none;">
      <div id="placeholderGenerador" class="img-placeholder" style="display:flex;">Aquí aparecerá tu imagen generada</div>
    </div>

    <!-- Prompt -->
    <textarea id="promptGenerador" placeholder="Describí la imagen que querés generar... (ej: un atardecer en la playa, estilo realista)" style="width:100%;height:75px;background:#001122;color:#00eaff;border:1px solid #006688;border-radius:10px;padding:10px 12px;font-size:14px;resize:vertical;box-sizing:border-box;outline:none;font-family:'Segoe UI',sans-serif;"></textarea>

    <!-- Botones -->
    <div style="display:flex;gap:10px;margin-top:14px;flex-wrap:wrap;align-items:center;">
      <button id="btnGenerarImagen" type="button" onclick="generarImagenIA()" style="padding:11px 24px;background:linear-gradient(135deg,#005577,#007799);color:#fff;border:none;border-radius:10px;font-size:15px;font-weight:600;cursor:pointer;box-shadow:0 0 14px #00eaff44;">✨ Generar con IA</button>
      <a id="btnDescargarGenerador" download="foschi_generada.png" style="display:none;padding:11px 18px;background:linear-gradient(135deg,#004400,#006600);color:#00ff88;border:1px solid #00ff8844;border-radius:10px;font-size:14px;font-weight:600;text-decoration:none;cursor:pointer;">⬇️ Descargar</a>
      <span id="generadorEstado" style="display:none;align-items:center;gap:8px;color:#00eaff88;font-size:13px;">
        <span style="display:inline-block;width:15px;height:15px;border:2px solid #00eaff33;border-top-color:#00eaff;border-radius:50%;animation:spinImg 0.7s linear infinite;"></span>
        Generando con IA...
      </span>
    </div>

  </div>
</div>

<!-- MODAL FLOTANTE GENERADOR DE PRESENTACIONES -->
<div id="generadorPresentacion" style="display:none;position:fixed;inset:0;z-index:9998;background:rgba(0,8,20,0.88);align-items:center;justify-content:center;padding:12px;box-sizing:border-box;">
  <div style="background:linear-gradient(135deg,#001a2e,#002a44);border:1px solid #00eaff55;border-radius:18px;box-shadow:0 0 40px #00eaff33;width:100%;max-width:640px;max-height:94vh;overflow-y:auto;padding:22px;box-sizing:border-box;">

    <!-- Título y botón cerrar -->
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:16px;">
      <span style="color:#00eaff;font-size:17px;font-weight:700;text-shadow:0 0 8px #00eaff;letter-spacing:1px;">🖥️ Crear Presentación con IA</span>
      <button onclick="cerrarGeneradorPresentacion()" style="background:transparent;border:1px solid #ff444466;color:#ff6666;border-radius:8px;padding:6px 16px;font-size:15px;cursor:pointer;">✕ Cerrar</button>
    </div>

    <div id="presInfoDoc" style="display:none;background:#00223a;border:1px solid #00eaff33;border-radius:10px;padding:10px 12px;color:#00eaff;font-size:13px;margin-bottom:12px;"></div>

    <label style="color:#00eaff99;font-size:12px;letter-spacing:1px;text-transform:uppercase;display:block;margin-bottom:6px;">Título (opcional)</label>
    <input id="presTitulo" type="text" placeholder="Ej: Plan de marketing 2026" style="width:100%;background:#001122;color:#00eaff;border:1px solid #006688;border-radius:10px;padding:10px 12px;font-size:14px;box-sizing:border-box;outline:none;font-family:'Segoe UI',sans-serif;margin-bottom:12px;">

    <label style="color:#00eaff99;font-size:12px;letter-spacing:1px;text-transform:uppercase;display:block;margin-bottom:6px;">Tema / qué querés que incluya</label>
    <textarea id="presTema" placeholder="Describí el tema, el público y los puntos clave que querés tratar... (si subiste un documento, esto se usa como enfoque opcional)" style="width:100%;height:90px;background:#001122;color:#00eaff;border:1px solid #006688;border-radius:10px;padding:10px 12px;font-size:14px;resize:vertical;box-sizing:border-box;outline:none;font-family:'Segoe UI',sans-serif;margin-bottom:12px;"></textarea>

    <div style="display:flex;gap:12px;flex-wrap:wrap;margin-bottom:12px;">
      <div style="flex:1;min-width:140px;">
        <label style="color:#00eaff99;font-size:12px;letter-spacing:1px;text-transform:uppercase;display:block;margin-bottom:6px;">Cantidad de diapositivas</label>
        <input id="presNumSlides" type="number" min="3" max="20" value="8" style="width:100%;background:#001122;color:#00eaff;border:1px solid #006688;border-radius:10px;padding:10px 12px;font-size:14px;box-sizing:border-box;outline:none;font-family:'Segoe UI',sans-serif;">
      </div>
      <div style="flex:1;min-width:140px;display:flex;align-items:center;">
        <label style="color:#00eaff;font-size:14px;display:flex;align-items:center;gap:8px;cursor:pointer;margin-top:18px;">
          <input id="presIncluirImagenes" type="checkbox" checked style="width:18px;height:18px;cursor:pointer;">
          🎨 Generar imágenes con IA
        </label>
      </div>
    </div>

    <label style="color:#00eaff99;font-size:12px;letter-spacing:1px;text-transform:uppercase;display:block;margin-bottom:6px;">Videos cortos (opcional, .mp4)</label>
    <input id="presVideos" type="file" accept="video/mp4,video/*" multiple style="width:100%;color:#00eaff;font-size:13px;margin-bottom:6px;">
    <div style="color:#00eaff66;font-size:12px;margin-bottom:14px;">Si subís videos, se insertan en las primeras diapositivas en lugar de las imágenes generadas.</div>

    <!-- Botones -->
    <div style="display:flex;gap:10px;margin-top:6px;flex-wrap:wrap;align-items:center;">
      <button id="btnGenerarPresentacion" type="button" onclick="generarPresentacionIA()" style="padding:11px 24px;background:linear-gradient(135deg,#005577,#007799);color:#fff;border:none;border-radius:10px;font-size:15px;font-weight:600;cursor:pointer;box-shadow:0 0 14px #00eaff44;">✨ Generar presentación</button>
      <span id="presEstado" style="display:none;align-items:center;gap:8px;color:#00eaff88;font-size:13px;">
        <span style="display:inline-block;width:15px;height:15px;border:2px solid #00eaff33;border-top-color:#00eaff;border-radius:50%;animation:spinImg 0.7s linear infinite;"></span>
        Generando presentación con IA... esto puede tardar 1-2 minutos, no cierres esta ventana.
      </span>
    </div>

  </div>
</div>
  #btnAplicarEdicion:disabled, #btnGenerarImagen:disabled { opacity:0.5; cursor:not-allowed; }
  .img-box{
    width:100%;
    height:280px;
    object-fit:contain;
    border-radius:10px;
    border:1px solid #00eaff22;
    background:#001122;
  }
  .img-placeholder{
    align-items:center;
    justify-content:center;
    height:280px;
    border-radius:10px;
    border:2px dashed #00eaff22;
    background:#001122;
    color:#00eaff33;
    font-size:13px;
    text-align:center;
    padding:0 12px;
    box-sizing:border-box;
  }
</style>
<div id="voiceWave" style="
position:fixed;
bottom:110px;
left:50%;
transform:translateX(-50%);
display:none;
align-items:center;
gap:4px;
z-index:999;
">
<div class="wave"></div>
<div class="wave"></div>
<div class="wave"></div>
<div class="wave"></div>
<div class="wave"></div>
<button onclick="detenerConversacion()" title="Salir del modo conversación" style="margin-left:10px;background:#ff0000;color:#fff;border:none;border-radius:50%;width:28px;height:28px;font-size:14px;font-weight:bold;line-height:1;cursor:pointer;display:flex;align-items:center;justify-content:center;box-shadow:0 0 10px #ff000066;">✕</button>
</div>

<div id="dictadoEstado" style="
 position:fixed;
 bottom:70px;
 right:10px;
 background:#ff0000;
 color:white;
 padding:6px 10px;
 border-radius:6px;
 display:none;
 font-weight:bold;
 z-index:999;">
🎤 Dictado activo
</div>
<!-- MENÚ ADJUNTOS (overlay fuera de la barra) -->
<div id="adjuntos_menu" aria-hidden="true">
  <button onclick="toggleModoConversacion()">🧠 Modo conversación</button>
  <button onclick="checkPremium('audio')">🎵 Audio (mp3/wav) a Texto</button>
  <button onclick="checkPremium('doc')">📄 Analizar Documento</button>
  <button onclick="abrirDictadoDesdeMenu()">🎤 Dictado</button>
  <button onclick="checkPremium('ocr')">
📷 Imagen a Word
</button>

<button onclick="checkPremium('editar_imagen')">
🖼️ Editar Imagen
</button>

<button onclick="abrirGeneradorImagen()">
🎨 Generar Imagen
</button>

<button onclick="abrirGeneradorPresentacion()">
🖥️ Crear Presentación
</button>

<button onclick="checkPremiumIngles()" style="border-top:1px solid #00cc6633;margin-top:4px;padding-top:10px;">
📚 Profesor de Inglés
</button>
</div>

<!-- INPUTS OCULTOS -->
<input id="audioInput" class="hidden_file_input" type="file" accept=".mp3,audio/*,.wav" />
<input
  id="imagenInput"
  class="hidden_file_input"
  type="file"
  accept="image/*"
/>
<input id="archivo_pdf_word" class="hidden_file_input" type="file" accept=".pdf,.docx" />

<!-- BARRA DE ENTRADA -->
<div id="inputBar">

  <!-- Clip — visible cuando NO está dictando -->
  <div id="clipBtn" title="Adjuntar" onclick="toggleAdjuntosMenu()">
    📎
    {% if not premium %}
    <span id="adjuntosHint">👉 Funciones Premium acá</span>
    {% endif %}
  </div>

  <!-- Input de texto -->
  <input type="text" id="mensaje" placeholder="Escribí tu mensaje o hablá" />

  <!-- Botones normales — se ocultan durante dictado -->
  <div id="botonesNormales" style="display:flex;gap:6px;align-items:center;">
    <button onclick="checkDailyLimit()">Enviar</button>
    <button onclick="hablar()">🎤 Hablar</button>
  </div>

  <!-- Botones dictado — solo visibles durante dictado -->
  <div id="botonesDictado" style="display:none;gap:6px;align-items:center;">
    <button id="dictadoBtn" onclick="toggleDictado()">⏸️ Pausar</button>
    <button id="finalizarDictadoBtn"
      onclick="finalizarDictadoManual()"
      style="background:#ff0033;color:white;border:1px solid #ff5577;">
      ✅ Finalizar
    </button>
    <button onclick="cancelarDictado()"
      style="background:#333;color:#ccc;border:1px solid #555;">
      ❌ Cancelar
    </button>
  </div>

</div>

<script>
// --- Variables y funciones generales ---
let usuario_id="{{usuario_id}}";
let documentoActual = null;
let textoDocumento = "";
let vozActiva=true,audioActual=null,mensajeActual=null;
let modoConversacion = false;
let escuchandoContinuo = false;
let recognitionConversacion = null;
let silencioTimer=null;
let modoImagen = "";
let imagenActualArchivo = null;
let imagenActualBase64 = null;

function toggleModoConversacion(){

  if(!isPremium && !isSuper){
    alert("🔒 El modo conversación es Premium");
    return;
  }

  if(!modoConversacion){
    iniciarConversacion();
  }else{
    detenerConversacion();
  }

}

function iniciarConversacion(){

  if(!('webkitSpeechRecognition' in window)){
    alert("Tu navegador no soporta conversación por voz");
    return;
  }

  document.getElementById("voiceWave").style.display="flex";


  const Rec = window.SpeechRecognition || window.webkitSpeechRecognition;

  recognitionConversacion = new Rec();

  recognitionConversacion.lang = "es-AR";
  recognitionConversacion.continuous = true;
  recognitionConversacion.interimResults = false;

  modoConversacion = true;
  escuchandoContinuo = true;

  agregar("🧠 Modo conversación activado","ai");

recognitionConversacion.onresult = function(event){

  let texto = event.results[event.results.length-1][0].transcript;
  if(texto.length < 3) return;

  if(texto.trim() === "") return;

  if(audioActual){
  audioActual.pause();
  audioActual.currentTime = 0;
}
  document.getElementById("mensaje").value = texto;

  enviar();

  // 👇 DETECTOR DE SILENCIO
  clearTimeout(silencioTimer);

  silencioTimer = setTimeout(()=>{
     console.log("Silencio detectado");
  },2000);

};

  recognitionConversacion.onend = function(){

  if(modoConversacion && !audioActual){
    safeStartRecognition();
  }

};

  safeStartRecognition();

}

// Evita "InvalidStateError: recognition has already started"
// cuando varios eventos intentan reiniciar el reconocimiento al mismo tiempo
function safeStartRecognition(){
  if(!recognitionConversacion) return;
  try{
    recognitionConversacion.start();
  }catch(e){
    // ya estaba escuchando, no hacemos nada
  }
}

function detenerConversacion(){

  modoConversacion = false;
  escuchandoContinuo = false;

  if(recognitionConversacion){
    recognitionConversacion.stop();
    recognitionConversacion = null;
  }

  document.getElementById("voiceWave").style.display = "none";

  agregar("🛑 Modo conversación desactivado","ai");

}

let MAX_NO_PREMIUM = 5;
let isPremium = {{ 'true' if premium else 'false' }};
const hoy = new Date().toISOString().slice(0,10); // YYYY-MM-DD
let preguntasHoy = parseInt(
  localStorage.getItem("preguntasHoy_" + hoy) || "0"
);

let isSuper = {{ 'true' if is_super else 'false' }};
let rolUsuario = "{{ rol or '' }}";
let nivelUsuario = {{ nivel or 0 }};

function logoClick(){ alert("FOSCHI NUNCA MUERE, TRASCIENDE..."); }
function toggleVoz(estado=null){ vozActiva=estado!==null?estado:!vozActiva; document.getElementById("vozBtn").textContent=vozActiva?"🔊 Voz activada":"🔇 Silenciada"; }
function detenerVoz(){ if(audioActual){ audioActual.pause(); audioActual.currentTime=0; audioActual.src=""; audioActual.load(); audioActual=null; if(mensajeActual) mensajeActual.classList.remove("playing"); mensajeActual=null; } }

function agregar(msg,cls,imagenes=[]){
  let c=document.getElementById("chat"),div=document.createElement("div");
  div.className="message "+cls; div.innerHTML=msg;
  c.appendChild(div);
  setTimeout(()=>div.classList.add("show"),50);
  imagenes.forEach(url=>{ let img=document.createElement("img"); img.src=url; div.appendChild(img); });
  c.scroll({top:c.scrollHeight,behavior:"smooth"});
  if(cls==="ai") hablarTexto(msg,div);
  return div;
}

function checkDailyLimit(){
  if(!isPremium && !isSuper && preguntasHoy >= MAX_NO_PREMIUM){
    alert(`⚠️ Has alcanzado el límite de ${MAX_NO_PREMIUM} preguntas diarias. Pasá a Premium para más.`);
    return;
  }

  enviar();

  if(!isPremium){
    preguntasHoy++;
    localStorage.setItem("preguntasHoy_" + hoy, preguntasHoy);
  }
}

function enviar(){
  let msg=document.getElementById("mensaje").value.trim(); 
  if(!msg) return;

  agregar(msg,"user"); 

  document.getElementById("mensaje").value="";
    // ============================
  // ❌ SALIR DEL MODO DOCUMENTO
  // ============================

  if(
    msg.toLowerCase().includes("salir del documento") ||
    msg.toLowerCase().includes("cerrar archivo") ||
    msg.toLowerCase().includes("modo normal")
  ){

    salirModoDocumento();

    return;
  }

  fetch("/preguntar",{
    method:"POST",
    headers:{
      "Content-Type":"application/json"
    },
    body:JSON.stringify({
      mensaje: msg,
      usuario_id: usuario_id,
      doc_id: documentoActual,
      preguntar_doc: modoPreguntasDocumento
    })
  })

  .then(r=>r.json())

  .then(data=>{

    agregar(data.texto,"ai",data.imagenes);

    if(data.borrar_historial){
      document.getElementById("chat").innerHTML="";
    }

  })

  .catch(e=>{

    agregar("Error al comunicarse con el servidor.","ai");

    console.error(e);

  });
}

document.getElementById("mensaje").addEventListener("keydown",e=>{ if(e.key==="Enter"){ e.preventDefault(); checkDailyLimit(); } });

// Saca emojis del texto para que la voz no los lea (ej: "manos en oración", "diamante", "mano saludando")
function quitarEmojisParaVoz(texto){
  return texto
    .replace(/[\\u{1F1E6}-\\u{1F1FF}\\u{1F300}-\\u{1FAFF}\\u{2190}-\\u{21FF}\\u{2300}-\\u{27BF}\\u{2B00}-\\u{2BFF}\\u{FE0F}\\u{200D}\\u{20E3}]/gu, "")
    .replace(/\s{2,}/g, " ")
    .trim();
}

function hablarTexto(texto, div=null){

  if(!vozActiva) return;

  // 🛑 detener escucha solo si está en modo conversación
  if(modoConversacion && recognitionConversacion){
    recognitionConversacion.stop();
  }

  detenerVoz();

  if(mensajeActual) mensajeActual.classList.remove("playing");
  if(div) div.classList.add("playing");
  mensajeActual = div;

  const textoParaVoz = quitarEmojisParaVoz(texto);
  audioActual = new Audio("/tts?texto=" + encodeURIComponent(textoParaVoz));
  audioActual.playbackRate = 1.25;

  audioActual.onended = () => {

    if(mensajeActual) mensajeActual.classList.remove("playing");
    mensajeActual = null;

    // 🎤 volver a escuchar cuando termina
    if(modoConversacion && recognitionConversacion){
      safeStartRecognition();
    }

  };

  audioActual.play();
}

function togglePremiumMenu(){
  if(isPremium) return;

  const menu = document.getElementById("premiumMenu");
  if(!menu) return;

  menu.style.display = (menu.style.display === "block") ? "none" : "block";

  if(menu.style.display === "block"){
    setTimeout(() => {
      window.addEventListener("click", closePremiumMenuOnClickOutside);
    }, 50);
  }
}

function closePremiumMenuOnClickOutside(e){
  const menu = document.getElementById("premiumMenu");
  const btn = document.getElementById("premiumBtn");
  if(!menu.contains(e.target) && !btn.contains(e.target)){
    menu.style.display="none";
    window.removeEventListener('click', closePremiumMenuOnClickOutside);
  }
}

function irPremium(tipo){
  fetch(`/premium?tipo=${tipo}`)
    .then(r => {
      if(r.status === 401){
        openAuth();
        return null;
      }
      return r.json();
    })
    .then(data => {
      if(!data) return;
      if(data.init_point){
        window.location.href = data.init_point;
      }
    })
    .catch(err => console.error(err));
}

function checkPremiumIngles() {
  if (!isPremium && !isSuper) {
    alert("⚠️ El Profesor de Inglés es una función Premium. Pasá a Premium para usarla.");
    return;
  }
  // Cerrar el menú de adjuntos
  const menu = document.getElementById("adjuntos_menu");
  if (menu) { menu.style.display = "none"; menu.setAttribute("aria-hidden","true"); }
  abrirProfesorIngles();
}

function checkPremium(tipo){

  if(!isPremium){
    alert("⚠️ Esta función requiere Premium. Pasá a Premium para usarla.");
    return;
  }

  if(tipo === "ocr"){
    modoImagen = "ocr";
    document.getElementById("imagenInput").click();
    return;
  }

  if(tipo === "editar_imagen"){

  modoImagen = "editar";

  document.getElementById(
    "imagenInput"
  ).click();

  return;
}

  if(tipo === "audio"){
    document.getElementById("audioInput").click();
    return;
  }

  if(tipo === "doc"){
    document.getElementById("archivo_pdf_word").click();
    return;
  }

}

let adjuntosInactivityTimer = null;

function iniciarTemporizadorInactividadAdjuntos(){
  // Se reinicia cada vez que el mouse se mueve mientras el menú está abierto
  document.addEventListener("mousemove", resetAdjuntosInactivityTimer);
  resetAdjuntosInactivityTimer();
}

function detenerTemporizadorInactividadAdjuntos(){
  clearTimeout(adjuntosInactivityTimer);
  adjuntosInactivityTimer = null;
  document.removeEventListener("mousemove", resetAdjuntosInactivityTimer);
}

function resetAdjuntosInactivityTimer(){
  clearTimeout(adjuntosInactivityTimer);
  adjuntosInactivityTimer = setTimeout(cerrarAdjuntosMenu, 6000); // 6s sin mover el mouse
}

function cerrarAdjuntosMenu(){
  const m = document.getElementById("adjuntos_menu");
  if(m) m.style.display = "none";
  window.removeEventListener('click', closeMenuOnClickOutside);
  detenerTemporizadorInactividadAdjuntos();
}

function toggleAdjuntosMenu(){
  // Una vez que el usuario abre el menú, ya no necesita el aviso
  const hint = document.getElementById("adjuntosHint");
  if(hint) hint.style.display = "none";

  const m = document.getElementById("adjuntos_menu");
  const seAbre = m.style.display !== "block";
  m.style.display = seAbre ? "block" : "none";

  if(seAbre){
    m.querySelectorAll("button").forEach(b => b.classList.remove("opcion-seleccionada"));
    setTimeout(()=>window.addEventListener('click', closeMenuOnClickOutside),50);
    iniciarTemporizadorInactividadAdjuntos();
  } else {
    detenerTemporizadorInactividadAdjuntos();
  }
}

// Cierra el menú de adjuntos automáticamente si el mouse se va de encima
document.getElementById("adjuntos_menu").addEventListener("mouseleave", cerrarAdjuntosMenu);

// Ilumina la opción elegida y luego cierra el menú
document.getElementById("adjuntos_menu").addEventListener("click", function(e){
  const btn = e.target.closest("button");
  if(!btn) return;

  document.querySelectorAll("#adjuntos_menu button").forEach(b => b.classList.remove("opcion-seleccionada"));
  btn.classList.add("opcion-seleccionada");

  setTimeout(cerrarAdjuntosMenu, 180);
});
function closeMenuOnClickOutside(e){
  const menu = document.getElementById("adjuntos_menu");
  const clip = document.getElementById("clipBtn");
  if(!menu.contains(e.target) && !clip.contains(e.target)){
    menu.style.display="none";
    window.removeEventListener('click', closeMenuOnClickOutside);
    detenerTemporizadorInactividadAdjuntos();
  }
}
function abrirDictadoDesdeMenu(){
  toggleDictado();
}

function verHistorial(){
  fetch("/historial/"+usuario_id).then(r=>r.json()).then(data=>{
    document.getElementById("chat").innerHTML="";
    if(data.length===0){agregar("No hay historial todavía.","ai");return;}
    data.slice(-20).forEach(e=>{ agregar(`<small>${e.fecha}</small><br>${e.usuario}`,"user"); agregar(`<small>${e.fecha}</small><br>${e.foschi}`,"ai"); });
  });
}

function borrarPantalla(){
    detenerVoz(); 
    document.getElementById("chat").innerHTML = ""; 
}

function openAuth() {
  document.getElementById("authModal").style.display = "block";
}

function closeAuth() {
  document.getElementById("authModal").style.display = "none";
}

async function login() {
  const email = authEmail.value;
  const password = authPassword.value;

  const r = await fetch("/auth/login", {
    method: "POST",
    headers: {"Content-Type":"application/json"},
    body: JSON.stringify({email, password})
  });

  const j = await r.json();
  if (j.ok) location.reload();
  else authMsg.innerText = j.msg;
}

async function register() {
  const email = authEmail.value;
  const password = authPassword.value;

  const r = await fetch("/auth/register", {
    method: "POST",
    headers: {"Content-Type":"application/json"},
    body: JSON.stringify({email, password})
  });

  const j = await r.json();
  if (j.ok) location.reload();
  else authMsg.innerText = j.msg;
}

function hablar(){
  if('webkitSpeechRecognition' in window || 'SpeechRecognition' in window){

    const Rec = window.SpeechRecognition || window.webkitSpeechRecognition;
    const recognition = new Rec();

    recognition.lang='es-AR';
    recognition.continuous=false;
    recognition.interimResults=false;

    recognition.onresult = function(event){

      let textoReconocido = event.results[0][0].transcript;
let txt = textoReconocido.toLowerCase();

// 👉 ACTIVAR DICTADO
if(txt.includes("activar dictado")){
  iniciarDictado();
  return;
}

// 👉 DESACTIVAR DICTADO
if(txt.includes("desactivar dictado")){
  detenerDictado();
  return;
}

// ============================
// BORRAR DESDE / HASTA
// ============================

let borrarMatch = txt.match(/borrar desde (.+) hasta (.+)/i);

if(borrarMatch){

  let desde = borrarMatch[1].trim();
  let hasta = borrarMatch[2].trim();

  let regex = new RegExp(
    desde + "[\\s\\S]*?" + hasta,
    "i"
  );

  let caja = document.getElementById("mensaje");

  caja.value = caja.value.replace(regex, "");

  alert("Texto eliminado");

  return;
}

// ============================
// CORREGIR PALABRA
// ============================

let corregirMatch = txt.match(
  /corregir (.+) por (.+)/i
);

if(corregirMatch){

  let original = corregirMatch[1].trim();
  let nuevo = corregirMatch[2].trim();

  let caja = document.getElementById("mensaje");

  let regex = new RegExp(original, "gi");

  caja.value = caja.value.replace(regex, nuevo);

  alert("Texto corregido");

  return;
}

// 👉 comportamiento normal
let caja = document.getElementById("mensaje");

caja.value += " " + txt;

checkDailyLimit();
    };

    recognition.onerror = function(e){
      console.log(e);
      alert("Error reconocimiento de voz: " + e.error);
    }

    recognition.start();

  }else{
    alert("Tu navegador no soporta reconocimiento de voz.");
  }
}

function chequearRecordatorios(){
  fetch("/avisos",{ method:"POST", headers:{"Content-Type":"application/json"}, body:JSON.stringify({usuario_id}) })
  .then(r=>r.json()).then(data=>{ if(Array.isArray(data) && data.length>0){ data.forEach(r=>{ agregar(`⏰ Tenés un recordatorio: ${r.motivo||"(sin motivo)"}`,"ai"); }); } }).catch(e=>console.error(e));
}
setInterval(chequearRecordatorios,10000);

/* --- SALUDO INICIAL --- */
window.onload = function() {
    let textoSaludo = isPremium
      ? "🙏 ¡Gracias por ser parte de Foschi IA Premium! 💎 Ya tenés todas las funciones desbloqueadas."
      : "👋 ¡Hola! Bienvenido a Foschi IA";

    agregar(textoSaludo,"ai");
};
/* =============================== */
/* === BOTÓN MODO DAY / NIGHT ==== */
/* =============================== */
function toggleDayNight(){
  const body = document.body;
  body.classList.toggle("day");

  const btn = document.getElementById("dayNightBtn");
  if(btn){
    btn.textContent = body.classList.contains("day") ? "☀️" : "🌙";
  }
}
// =====================
// 🎤 DICTADO PRO
// =====================

let dictadoActivo = false;
let dictadoPausado = false;
let reconocimiento = null;
let textoDictado = localStorage.getItem("dictado_guardado") || "";
let ultimoTexto = "";
let reinicioDictado = false;

// Crea el panel de previsualización una sola vez
function getPanelDictado(){
  let p = document.getElementById("dictadoPanel");
  if(!p){
    p = document.createElement("div");
    p.id = "dictadoPanel";
    document.body.appendChild(p);
  }
  return p;
}

// Refresca lo que se ve en el panel:
// texto confirmado (blanco) + lo que se está diciendo ahora (gris)
function actualizarPanel(parcial){
  const panel = getPanelDictado();
  const enCurso = parcial
    ? '<span style="color:#888">' + parcial + '</span>'
    : "";
  panel.innerHTML = textoDictado + enCurso;
  panel.scrollTop = panel.scrollHeight;
}

function _resetUI(){
  document.getElementById("dictadoEstado").style.display = "none";
  document.getElementById("clipBtn").style.display = "flex";
  document.getElementById("botonesNormales").style.display = "flex";
  document.getElementById("botonesDictado").style.display = "none";
  document.getElementById("mensaje").value = "";
  document.getElementById("mensaje").placeholder = "Escribí tu mensaje o hablá";
  const panel = getPanelDictado();
  panel.style.display = "none";
  panel.innerHTML = "";
}

function toggleDictado(){
  if(!isPremium && !isSuper){
    alert("🔒 Esta función es solo Premium");
    return;
  }
  if(!dictadoActivo){ iniciarDictado(); return; }
  if(dictadoActivo && !dictadoPausado){ pausarDictado(); return; }
  if(dictadoActivo && dictadoPausado){ continuarDictado(); }
}

function iniciarDictado(){

  const SpeechRecognition =
    window.SpeechRecognition || window.webkitSpeechRecognition;

  if(!SpeechRecognition){
    alert("Tu navegador no soporta dictado");
    return;
  }

  reconocimiento = new SpeechRecognition();
  reconocimiento.lang = "es-AR";
  reconocimiento.continuous = true;
  reconocimiento.interimResults = true;
  reconocimiento.maxAlternatives = 1;

  dictadoActivo = true;
  dictadoPausado = false;
  reinicioDictado = true;
  if(!textoDictado) textoDictado = "";
  ultimoTexto = "";

  // UI: estado visible, clip+botones normales ocultos, botones dictado visibles
  document.getElementById("dictadoEstado").style.display = "block";
  document.getElementById("dictadoEstado").innerText = "🎤 Dictando...";
  document.getElementById("clipBtn").style.display = "none";
  document.getElementById("botonesNormales").style.display = "none";
  document.getElementById("botonesDictado").style.display = "flex";
  document.getElementById("adjuntos_menu").style.display = "none";
  document.getElementById("dictadoBtn").classList.add("activo");
  document.getElementById("dictadoBtn").innerText = "⏸️ Pausar";

  // Mostrar panel y limpiar input
  const panel = getPanelDictado();
  panel.style.display = "block";
  actualizarPanel("");
  document.getElementById("mensaje").value = "";
  document.getElementById("mensaje").placeholder = "🎤 Dictando...";

 reconocimiento.onresult = function(event){

    let parcial = "";

    for(let i = event.resultIndex; i < event.results.length; i++){

      let trans = event.results[i][0].transcript.trim();

      if(!trans) continue;

      let txt = trans.toLowerCase();

      // ============================
      // 🗑️ BORRAR PALABRAS
      // ============================

      // borrar eso
      if(
        txt.includes("borrar eso") ||
        txt.includes("borrar esa palabra")
      ){

        let palabras = textoDictado.trim().split(" ");

        palabras.pop();

        textoDictado = palabras.join(" ");

        document.getElementById("mensaje").value = textoDictado;

        return;
      }

      // ============================
// 🗑️ BORRAR X PALABRAS
// ============================

if(
  txt.includes("borrar") &&
  txt.includes("palabras")
){

  let numero = 1;

  // números escritos
  if(txt.includes("dos")) numero = 2;
  else if(txt.includes("tres")) numero = 3;
  else if(txt.includes("cuatro")) numero = 4;
  else if(txt.includes("cinco")) numero = 5;
  else if(txt.includes("seis")) numero = 6;
  else if(txt.includes("siete")) numero = 7;
  else if(txt.includes("ocho")) numero = 8;
  else if(txt.includes("nueve")) numero = 9;
  else if(txt.includes("diez")) numero = 10;

  // números normales
  let match = txt.match(/\d+/);

  if(match){
    numero = parseInt(match[0]);
  }

  let palabras = textoDictado.trim().split(/\s+/);

  palabras.splice(-numero);

  textoDictado = palabras.join(" ");

  document.getElementById("mensaje").value = textoDictado;

  return;
}

      // ============================
      // 🎤 COMANDOS DE VOZ
      // ============================

      if(txt.includes("pausar dictado"))   {
        pausarDictado();
        return;
      }

      if(txt.includes("continuar dictado")){
        continuarDictado();
        return;
      }

      if(txt.includes("finalizar dictado")){
        finalizarDictado();
        return;
      }

      if(
        txt.includes("cancelar dictado") ||
        txt.includes("borrar dictado") ||
        txt.includes("borrar todo")
      ){
        cancelarDictado();
        return;
      }

      trans = mejorarTextoDictado(trans);

      if(event.results[i].isFinal){

        // Evitar duplicado del último fragmento final
        if(trans !== ultimoTexto){

          ultimoTexto = trans;

          textoDictado += trans + " ";

          localStorage.setItem(
            "dictado_guardado",
            textoDictado
          );
        }

      } else {

        parcial += trans + " ";
      }
    }

    // Refrescar panel con acumulado + parcial en gris
    actualizarPanel(parcial);

    // El input solo muestra conteo para no distraer
    const palabras = textoDictado.trim().split(/\s+/).filter(Boolean).length;
    document.getElementById("mensaje").placeholder =
      "🎤 " + palabras + " palabra" + (palabras !== 1 ? "s" : "") + " dictada" + (palabras !== 1 ? "s" : "");
  };

  reconocimiento.onerror = function(e){
    console.log("Error dictado:", e.error);
    if(e.error === "not-allowed"){
      alert("Micrófono bloqueado");
      finalizarDictado();
    }
  };

  reconocimiento.onend = function(){
    if(dictadoActivo && !dictadoPausado && reinicioDictado){
      try{ reconocimiento.start(); }catch(err){ console.log("Reinicio:", err); }
    }
  };

  reconocimiento.start();
}

function pausarDictado(){
  dictadoPausado = true;
  if(reconocimiento){ reinicioDictado = false; reconocimiento.stop(); }
  document.getElementById("dictadoEstado").innerText = "⏸️ Dictado pausado";
  document.getElementById("dictadoBtn").innerText = "▶️ Continuar";
}

function continuarDictado(){
  if(!dictadoActivo) return;
  dictadoPausado = false;
  reinicioDictado = true;
  try{ reconocimiento.start(); }catch(err){ console.log("Continuar:", err); }
  document.getElementById("dictadoEstado").innerText = "🎤 Dictando...";
  document.getElementById("dictadoBtn").innerText = "⏸️ Pausar";
}

function finalizarDictado(){
  dictadoActivo = false;
  dictadoPausado = false;
  reinicioDictado = false;

  if(reconocimiento){ reconocimiento.stop(); reconocimiento = null; }

  // ⚠️ Guardar el texto ANTES de limpiar las variables
  const textoFinal = textoDictado.trim();
  textoDictado = "";
  ultimoTexto = "";
  localStorage.removeItem("dictado_guardado");

  _resetUI();

  if(textoFinal.length > 0){
    descargarWordDictado(textoFinal);
  } else {
    alert("No hay texto para guardar");
  }
}

function finalizarDictadoManual(){
  if(!textoDictado.trim()){ alert("No hay texto para guardar"); return; }
  finalizarDictado();
}

function detenerDictado(){ finalizarDictado(); }

function cancelarDictado(){
  if(reconocimiento){ reconocimiento.stop(); }
  dictadoActivo = false;
  dictadoPausado = false;
  reinicioDictado = false;
  textoDictado = "";
  ultimoTexto = "";
  localStorage.removeItem("dictado_guardado");
  _resetUI();
  agregar("🗑️ Dictado cancelado", "ai");
}

// =====================
// MEJORADOR DE TEXTO
// =====================

function mejorarTextoDictado(texto){

  texto = texto.trim();

  // =========================
  // PUNTUACIÓN PROFESIONAL
  // =========================

  texto = texto.replace(/coma/gi, ",");
  texto = texto.replace(/\bdos puntos\b/gi, ":");

  texto = texto.replace(/punto y coma/gi, ";");
  texto = texto.replace(/punto aparte/gi, ".");
  texto = texto.replace(/punto seguido/gi, ".");
  texto = texto.replace(/punto/gi, ".");

  // preguntas
  texto = texto.replace(/\babrir pregunta\b/gi, "¿");
  texto = texto.replace(/\bcerrar pregunta\b/gi, "?");

  // admiración
  texto = texto.replace(/\babrir admiración\b/gi, "¡");
  texto = texto.replace(/\bcerrar admiración\b/gi, "!");

  // paréntesis
  texto = texto.replace(/\babrir paréntesis\b/gi, "(");
  texto = texto.replace(/\bcerrar paréntesis\b/gi, ")");

  // nuevo párrafo
  texto = texto.replace(/nuevo párrafo/gi, " ");

  // =========================
  // ESPACIOS
  // =========================

  // quitar espacios antes de signos
  texto = texto.replace(/\s+([.,;:!?])/g, "$1");

  // agregar espacio después de signos
  texto = texto.replace(/([.,;:!?])([A-Za-zÁÉÍÓÚáéíóúÑñ])/g, "$1 $2");

  // eliminar espacios dobles
  texto = texto.replace(/\s{2,}/g, " ");

  // =========================
  // MAYÚSCULAS
  // =========================

  if(
    textoDictado.trim() === "" ||
    /[.!?]\s*$/.test(textoDictado.trim())
  ){
    texto =
      texto.charAt(0).toUpperCase() +
      texto.slice(1);
  }

  // mayúscula después de punto
  texto = texto.replace(
    /([.!?]\s*)([a-záéíóúñ])/g,
    function(match, p1, p2){
      return p1 + p2.toUpperCase();
    }
  );

  return texto;
}

async function descargarWordDictado(texto){

  try{

    const r = await fetch("/dictado_word",{
      method:"POST",
      headers:{
        "Content-Type":"application/json"
      },
      body:JSON.stringify({
        texto:texto
      })
    });

    if(!r.ok){

      let errorTexto = await r.text();

      console.log(errorTexto);

      alert("Error generando Word");

      return;
    }

    const blob = await r.blob();

    const url =
      window.URL.createObjectURL(blob);

    const a = document.createElement("a");

    a.href = url;

    a.download =
      "dictado_foschi.docx";

    document.body.appendChild(a);

    a.click();

    a.remove();

    window.URL.revokeObjectURL(url);

    // ✅ LIMPIAR SOLO DESPUÉS
    localStorage.removeItem(
      "dictado_guardado"
    );

  }catch(err){

    console.log(err);

    alert("Error descargando Word");
  }
}

// ===============================
// 🎵 AUDIO A TEXTO (WORD)
// ===============================

document.getElementById("audioInput")
.addEventListener("change", async function(e){

  const file = e.target.files[0];
  if(!file) return;

  agregar("🎵 Transcribiendo audio, esperá un momento...","ai");

  let formData = new FormData();
  formData.append("audio", file);
  formData.append("usuario_id", usuario_id);

  try{

    const r = await fetch("/upload_audio",{
      method:"POST",
      body:formData
    });

    if(!r.ok){
      let txt = await r.text();
      agregar("❌ Error en transcripción: " + txt, "ai");
      e.target.value = "";
      return;
    }

    const blob = await r.blob();
    const url = window.URL.createObjectURL(blob);
    const a = document.createElement("a");
    a.href = url;
    a.download = file.name.replace(/\.[^.]+$/, "") + "_transcripcion.docx";
    document.body.appendChild(a);
    a.click();
    a.remove();
    window.URL.revokeObjectURL(url);

    agregar("✅ Transcripción lista. Se descargó el Word con el texto.", "ai");

  }catch(err){
    console.log(err);
    agregar("❌ Error procesando el audio.", "ai");
  }

  e.target.value = ""; // reset para poder subir el mismo archivo de nuevo

});

// ===============================
// 📷 IMAGEN A WORD
// ===============================

document.getElementById("imagenInput")
.addEventListener("change", async function(e){

  const file = e.target.files[0];

  if(!file) return;

  // 🖼️ EDITAR IMAGEN
  if(modoImagen === "editar"){

    imagenActualArchivo = file;

    let lector = new FileReader();

    lector.onload = function(ev){

      // Mostrar imagen original
      let prev = document.getElementById("previewImagen");
      prev.src = ev.target.result;
      prev.style.display = "block";
      let phOrig = document.getElementById("placeholderOrig");
      if(phOrig) phOrig.style.display = "none";

      // Resetear resultado
      let res = document.getElementById("resultadoImagen");
      res.src = ""; res.style.display = "none";
      let phRes = document.getElementById("placeholderResult");
      if(phRes) phRes.style.display = "flex";

      // Resetear controles
      document.getElementById("btnDescargarEdicion").style.display = "none";
      document.getElementById("editorEstado").style.display = "none";
      document.getElementById("btnAplicarEdicion").disabled = false;
      document.getElementById("promptImagen").value = "";

      // Abrir modal flotante
      document.getElementById("editorImagen").style.display = "flex";
    };

    lector.readAsDataURL(file);

    return;
  }

  // 📷 IMAGEN A WORD
  let formData = new FormData();

  formData.append("imagen", file);

  agregar("📷 Procesando imagen...", "ai");

  try{

    const r = await fetch(
      "/imagen_a_word",
      {
        method:"POST",
        body:formData
      }
    );

    if(!r.ok){

      let txt = await r.text();

      agregar(
        "❌ Error: " + txt,
        "ai"
      );

      e.target.value = "";

      return;
    }

    const blob = await r.blob();

    const url =
      window.URL.createObjectURL(blob);

    const a =
      document.createElement("a");

    a.href = url;

    a.download =
      "imagen_extraida.docx";

    document.body.appendChild(a);

    a.click();

    a.remove();

    window.URL.revokeObjectURL(url);

    agregar(
      "✅ Word generado correctamente.",
      "ai"
    );

  }catch(err){

    console.log(err);

    agregar(
      "❌ Error procesando la imagen.",
      "ai"
    );
  }

  e.target.value = "";

});

// ===============================
// 📄 SUBIR PDF / WORD
// ===============================

document.getElementById("archivo_pdf_word")
.addEventListener("change", async function(e){

  const file = e.target.files[0];

  if(!file) return;

  agregar("📄 Subiendo documento...","ai");

  let formData = new FormData();

  formData.append("archivo", file);
  formData.append("usuario_id", usuario_id);

  try{

    const r = await fetch("/upload_doc",{
      method:"POST",
      body:formData
    });

    if(!r.ok){
      let txt = await r.text();
      agregar("❌ Error: " + txt,"ai");
      return;
    }

    const data = await r.json();

    documentoActual = data.doc_id;
    textoDocumento = data.snippet || "";

    agregar(
      `✅ Documento cargado: ${data.name}
      <br><br>
      📌 Elegí una opción:
      <br><br>
      <button onclick="resumirDocumento('breve')">📄 Resumen breve</button>
      <button onclick="resumirDocumento('normal')">📘 Resumen normal</button>
      <button onclick="resumirDocumento('profundo')">📚 Resumen profundo</button>
      <button onclick="activarPreguntasDocumento()">❓ Preguntar al documento</button>
      `,
      "ai"
    );

  }catch(err){

    console.log(err);

    agregar(
      "❌ Error subiendo documento",
      "ai"
    );
  }

});

// ===============================
// 📘 RESUMIR DOCUMENTO
// ===============================

async function resumirDocumento(tipo){

  if(!documentoActual){
    agregar("❌ No hay documento cargado","ai");
    return;
  }

  agregar("🧠 Generando resumen...","ai");

  try{

    const r = await fetch("/resumir_doc",{
      method:"POST",
      headers:{
        "Content-Type":"application/json"
      },
      body:JSON.stringify({
        doc_id: documentoActual,
        modo: tipo,
        usuario_id: usuario_id
      })
    });

    if(!r.ok){

      let txt = await r.text();

      agregar("❌ Error: " + txt,"ai");

      return;
    }

    const blob = await r.blob();

    const url = window.URL.createObjectURL(blob);

    const a = document.createElement("a");

    a.href = url;

    a.download = "resumen_foschi.docx";

    document.body.appendChild(a);

    a.click();

    a.remove();

    agregar("✅ Resumen generado","ai");

  }catch(err){

    console.log(err);

    agregar("❌ Error generando resumen","ai");
  }
}

// ===============================
// ❓ PREGUNTAS SOBRE DOCUMENTO
// ===============================

let modoPreguntasDocumento = false;

function activarPreguntasDocumento(){

  modoPreguntasDocumento = true;

  agregar(
    "📄 Modo documento ACTIVADO. Ahora podés hacer preguntas sobre el archivo.",
    "ai"
  );

  // ============================
  // CREAR BOTÓN FLOTANTE
  // ============================

  if(!document.getElementById("btnSalirDocumento")){

    let btn = document.createElement("button");

    btn.id = "btnSalirDocumento";

    btn.innerHTML = "❌ Salir del documento";

    btn.onclick = salirModoDocumento;

    btn.style.position = "fixed";

    btn.style.bottom = "90px";

    btn.style.right = "20px";

    btn.style.zIndex = "9999";

    btn.style.padding = "12px 18px";

    btn.style.borderRadius = "14px";

    btn.style.border = "none";

    btn.style.cursor = "pointer";

    btn.style.fontWeight = "bold";

    btn.style.background = "#ff0033";

    btn.style.color = "white";

    btn.style.boxShadow = "0 0 15px rgba(255,0,0,0.5)";

    document.body.appendChild(btn);
  }
}

function salirModoDocumento(){

  modoPreguntasDocumento = false;

  documentoActual = null;

  textoDocumento = "";

  let btn = document.getElementById("btnSalirDocumento");

  if(btn){
    btn.remove();
  }

  agregar(
    "✅ Saliste del modo documento. Foschi IA volvió al modo normal.",
    "ai"
  );
}

async function editarImagenActual(){

    let prompt = document
        .getElementById("promptImagen")
        .value
        .trim();

    if(!prompt){
        alert("Escribí una modificación");
        return;
    }

    // UI: cargando
    let btnAplicar = document.getElementById("btnAplicarEdicion");
    let estado = document.getElementById("editorEstado");
    btnAplicar.disabled = true;
    estado.style.display = "flex";
    document.getElementById("btnDescargarEdicion").style.display = "none";

    // Ocultar resultado anterior
    let resImg = document.getElementById("resultadoImagen");
    let phRes  = document.getElementById("placeholderResult");
    resImg.style.display = "none";
    if(phRes) phRes.style.display = "flex";

    let formData = new FormData();
    formData.append("imagen", imagenActualArchivo);
    formData.append("prompt", prompt);

    try{

        const r = await fetch("/editar_imagen", { method:"POST", body:formData });
        const data = await r.json();

        estado.style.display = "none";
        btnAplicar.disabled = false;

        if(!data.ok){
            alert(data.error);
            return;
        }

        // Mostrar resultado en columna derecha
        let dataUrl = "data:image/png;base64," + data.imagen;
        resImg.src = dataUrl;
        resImg.style.display = "block";
        if(phRes) phRes.style.display = "none";

        // Habilitar descarga
        let btnDesc = document.getElementById("btnDescargarEdicion");
        btnDesc.href = dataUrl;
        btnDesc.style.display = "inline-block";

        agregar("✅ Imagen editada con IA", "ai");

    }catch(err){

        console.log(err);
        estado.style.display = "none";
        btnAplicar.disabled = false;
        agregar("❌ Error editando imagen", "ai");
    }
}
function cancelarEdicionImagen(){

    document.getElementById("editorImagen").style.display = "none";
    document.getElementById("promptImagen").value = "";

    let prev = document.getElementById("previewImagen");
    prev.src = "";
    prev.style.display = "none";
    let phOrig = document.getElementById("placeholderOrig");
    if(phOrig) phOrig.style.display = "flex";

    let res = document.getElementById("resultadoImagen");
    res.src = "";
    res.style.display = "none";
    let phRes = document.getElementById("placeholderResult");
    if(phRes) phRes.style.display = "flex";

    let btnDesc = document.getElementById("btnDescargarEdicion");
    if(btnDesc) btnDesc.style.display = "none";
    let estado = document.getElementById("editorEstado");
    if(estado) estado.style.display = "none";
    let btnAplicar = document.getElementById("btnAplicarEdicion");
    if(btnAplicar) btnAplicar.disabled = false;

    imagenActualArchivo = null;
    modoImagen = "";
}

// ===============================
// 🎨 GENERADOR DE IMAGEN IA
// ===============================

function abrirGeneradorImagen(){

    if(!isPremium){
        alert("⚠️ Esta función requiere Premium. Pasá a Premium para usarla.");
        return;
    }

    document.getElementById("promptGenerador").value = "";

    let res = document.getElementById("resultadoGenerador");
    res.src = "";
    res.style.display = "none";

    let ph = document.getElementById("placeholderGenerador");
    if(ph) ph.style.display = "flex";

    document.getElementById("btnDescargarGenerador").style.display = "none";
    document.getElementById("generadorEstado").style.display = "none";
    document.getElementById("btnGenerarImagen").disabled = false;

    document.getElementById("generadorImagen").style.display = "flex";
}

async function generarImagenIA(){

    let prompt = document
        .getElementById("promptGenerador")
        .value
        .trim();

    if(!prompt){
        alert("Describí la imagen que querés generar");
        return;
    }

    // UI: cargando
    let btnGenerar = document.getElementById("btnGenerarImagen");
    let estado = document.getElementById("generadorEstado");
    btnGenerar.disabled = true;
    estado.style.display = "flex";
    document.getElementById("btnDescargarGenerador").style.display = "none";

    // Ocultar resultado anterior
    let resImg = document.getElementById("resultadoGenerador");
    let phRes  = document.getElementById("placeholderGenerador");
    resImg.style.display = "none";
    if(phRes) phRes.style.display = "flex";

    try{

        const r = await fetch("/generar_imagen", {
            method: "POST",
            headers: { "Content-Type": "application/json" },
            body: JSON.stringify({ prompt: prompt })
        });
        const data = await r.json();

        estado.style.display = "none";
        btnGenerar.disabled = false;

        if(!data.ok){
            alert(data.error);
            return;
        }

        // Mostrar resultado
        let dataUrl = "data:image/png;base64," + data.imagen;
        resImg.src = dataUrl;
        resImg.style.display = "block";
        if(phRes) phRes.style.display = "none";

        // Habilitar descarga
        let btnDesc = document.getElementById("btnDescargarGenerador");
        btnDesc.href = dataUrl;
        btnDesc.style.display = "inline-block";

        agregar("✅ Imagen generada con IA", "ai");

    }catch(err){

        console.log(err);
        estado.style.display = "none";
        btnGenerar.disabled = false;
        agregar("❌ Error generando imagen", "ai");
    }
}

function cancelarGeneradorImagen(){

    document.getElementById("generadorImagen").style.display = "none";
    document.getElementById("promptGenerador").value = "";

    let res = document.getElementById("resultadoGenerador");
    res.src = "";
    res.style.display = "none";
    let ph = document.getElementById("placeholderGenerador");
    if(ph) ph.style.display = "flex";

    let btnDesc = document.getElementById("btnDescargarGenerador");
    if(btnDesc) btnDesc.style.display = "none";
    let estado = document.getElementById("generadorEstado");
    if(estado) estado.style.display = "none";
    let btnGenerar = document.getElementById("btnGenerarImagen");
    if(btnGenerar) btnGenerar.disabled = false;
}

// ===============================
// 🖥️ GENERADOR DE PRESENTACIONES (PPTX)
// ===============================

function abrirGeneradorPresentacion(temaPrefill){

    if(!isPremium){
        alert("⚠️ Esta función requiere Premium. Pasá a Premium para usarla.");
        return;
    }

    document.getElementById("presTitulo").value = "";
    document.getElementById("presTema").value = temaPrefill || "";
    document.getElementById("presNumSlides").value = 8;
    document.getElementById("presIncluirImagenes").checked = true;
    document.getElementById("presVideos").value = "";
    document.getElementById("presEstado").style.display = "none";
    document.getElementById("btnGenerarPresentacion").disabled = false;

    let infoDoc = document.getElementById("presInfoDoc");
    if(documentoActual){
        infoDoc.style.display = "block";
        infoDoc.innerHTML = "📄 Tenés un documento cargado: se usará como base de la presentación. Si además escribís un tema, lo usaremos como enfoque adicional.";
    }else{
        infoDoc.style.display = "none";
    }

    document.getElementById("generadorPresentacion").style.display = "flex";
}

function cerrarGeneradorPresentacion(){
    document.getElementById("generadorPresentacion").style.display = "none";
}

async function generarPresentacionIA(){

    let tema = document.getElementById("presTema").value.trim();
    let titulo = document.getElementById("presTitulo").value.trim();
    let numSlides = document.getElementById("presNumSlides").value || 8;
    let incluirImagenes = document.getElementById("presIncluirImagenes").checked;
    let videos = document.getElementById("presVideos").files;

    if(!tema && !documentoActual){
        alert("Escribí un tema/descripción o subí un documento (📄 Analizar Documento) antes de generar la presentación.");
        return;
    }

    let btn = document.getElementById("btnGenerarPresentacion");
    let estado = document.getElementById("presEstado");
    btn.disabled = true;
    estado.style.display = "flex";

    let formData = new FormData();
    formData.append("usuario_id", usuario_id);
    formData.append("tema", tema);
    formData.append("titulo", titulo);
    formData.append("num_slides", numSlides);
    formData.append("incluir_imagenes", incluirImagenes ? "true" : "false");

    if(documentoActual){
        formData.append("doc_id", documentoActual);
    }

    for(let i=0; i<videos.length; i++){
        formData.append("videos", videos[i]);
    }

    try{

        // 1) Iniciar el job de generación (responde rápido, sin esperar la IA)
        const r = await fetch("/generar_presentacion", {
            method: "POST",
            body: formData
        });

        let data = null;
        try{ data = await r.json(); }catch(e){}

        if(!r.ok || !data || !data.ok || !data.job_id){
            let errTxt = (data && data.error) ? data.error : "Error iniciando la generación de la presentación";
            estado.style.display = "none";
            btn.disabled = false;
            alert("❌ " + errTxt);
            return;
        }

        // 2) Consultar el estado periódicamente hasta que esté listo
        await esperarYDescargarPresentacion(data.job_id, btn, estado);

    }catch(err){

        console.log(err);
        estado.style.display = "none";
        btn.disabled = false;
        agregar("❌ Error generando la presentación", "ai");
    }
}

function esperarYDescargarPresentacion(jobId, btn, estado){
    return new Promise((resolve)=>{
        const intervalo = setInterval(async ()=>{
            try{
                const r = await fetch("/estado_presentacion/" + jobId);
                const data = await r.json();

                if(!data.ok){
                    clearInterval(intervalo);
                    estado.style.display = "none";
                    btn.disabled = false;
                    alert("❌ " + (data.error || "Error consultando el estado de la presentación"));
                    resolve();
                    return;
                }

                if(data.status === "listo"){
                    clearInterval(intervalo);

                    // Descargar el archivo ya generado
                    const a = document.createElement("a");
                    a.href = "/descargar_presentacion/" + jobId;
                    a.download = "presentacion_foschi.pptx";
                    document.body.appendChild(a);
                    a.click();
                    a.remove();

                    estado.style.display = "none";
                    btn.disabled = false;
                    cerrarGeneradorPresentacion();
                    agregar("✅ Presentación generada y descargada (.pptx)", "ai");
                    resolve();

                }else if(data.status === "error"){
                    clearInterval(intervalo);
                    estado.style.display = "none";
                    btn.disabled = false;
                    alert("❌ " + (data.error || "Error generando la presentación"));
                    resolve();
                }
                // si status === "procesando", seguimos esperando

            }catch(err){
                clearInterval(intervalo);
                estado.style.display = "none";
                btn.disabled = false;
                console.log(err);
                agregar("❌ Error consultando el estado de la presentación", "ai");
                resolve();
            }
        }, 4000);
    });
}

</script>

<div id="authModal" style="display:none; position:fixed; inset:0; background:rgba(0,0,0,.85); z-index:9999;">
  <div style="max-width:360px; margin:10% auto; background:#001d3d; padding:20px; border-radius:12px; color:#00eaff;">
    
    <h3>Ingresar</h3>

    <div class="form-premium">

      <input id="authEmail" type="email" placeholder="Email">

      <input id="authPassword" type="password" placeholder="Contraseña">

      <button onclick="login()">Ingresar</button>

      <button onclick="register()">Crear cuenta</button>

    </div>

    <p id="authMsg" style="margin-top:10px;"></p>


    <button onclick="closeAuth()" style="margin-top:10px;">Cerrar</button>

  </div>
</div>

<!-- ═══════════════════════════════════════════════════════ -->
<!-- 📚  MODAL PROFESOR DE INGLÉS                           -->
<!-- ═══════════════════════════════════════════════════════ -->
<div id="modalIngles" style="display:none;position:fixed;inset:0;z-index:9999;background:rgba(0,8,20,0.95);align-items:flex-start;justify-content:center;padding:10px;box-sizing:border-box;overflow-y:auto;">
  <div style="background:linear-gradient(160deg,#001a0d,#002a1a,#001a2e);border:1.5px solid #00cc6644;border-radius:20px;box-shadow:0 0 50px #00cc6622;width:100%;max-width:820px;margin:0 auto;padding:22px;box-sizing:border-box;">

    <!-- HEADER -->
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:18px;">
      <div>
        <span style="color:#00ff88;font-size:20px;font-weight:800;text-shadow:0 0 12px #00ff8888;">📚 Profesor de Inglés</span>
        <span id="inglesNivelBadge" style="margin-left:12px;background:#003d1a;border:1px solid #00cc66;color:#00ff88;border-radius:20px;padding:3px 12px;font-size:12px;font-weight:700;">A1</span>
      </div>
      <button onclick="cerrarIngles()" style="background:transparent;border:1px solid #ff444466;color:#ff6666;border-radius:8px;padding:6px 16px;cursor:pointer;">✕ Cerrar</button>
    </div>

    <!-- TABS -->
    <div style="display:flex;gap:8px;margin-bottom:18px;flex-wrap:wrap;">
      <button class="ingles-tab active" onclick="inglesSetTab('conversacion')" id="tab-conversacion">💬 Conversación</button>
      <button class="ingles-tab" onclick="inglesSetTab('leccion')" id="tab-leccion">🎓 Lección</button>
      <button class="ingles-tab" onclick="inglesSetTab('escenario')" id="tab-escenario">🌍 Escenario real</button>
      <button class="ingles-tab" onclick="inglesSetTab('progreso')" id="tab-progreso">📊 Mi progreso</button>
    </div>

    <!-- ── CONVERSACIÓN ── -->
    <div id="panel-conversacion" class="ingles-panel">
      <div style="margin-bottom:10px;color:#00cc66;font-size:13px;">Hablá con tu profesor en inglés. Corrige tus errores y explica en español. 🇬🇧</div>
      <div id="inglesChat" style="background:#000d06;border:1px solid #00441a;border-radius:12px;padding:14px;height:320px;overflow-y:auto;margin-bottom:12px;display:flex;flex-direction:column;gap:10px;"></div>
      <div style="display:flex;gap:8px;align-items:flex-end;">
        <textarea id="inglesInput" placeholder="Write in English..." style="flex:1;background:#001a0d;color:#e6ffe6;border:1px solid #00cc6644;border-radius:10px;padding:10px 12px;font-size:14px;resize:none;height:52px;outline:none;font-family:'Segoe UI',sans-serif;" onkeydown="if(event.key==='Enter'&&!event.shiftKey){event.preventDefault();inglesEnviar();}"></textarea>
        <button onclick="inglesEnviar()" style="background:linear-gradient(135deg,#003d1a,#006633);color:#00ff88;border:1px solid #00cc66;border-radius:10px;padding:12px 18px;cursor:pointer;font-size:18px;">➤</button>
        <button onclick="inglesHablar()" id="inglesVozBtn" title="Hablar en inglés" style="background:#001a0d;color:#00cc66;border:1px solid #00cc6644;border-radius:10px;padding:12px 14px;cursor:pointer;font-size:18px;">🎤</button>
      </div>
      <div style="margin-top:6px;text-align:right;">
        <button onclick="document.getElementById('inglesChat').innerHTML=''" style="background:transparent;border:none;color:#006633;font-size:12px;cursor:pointer;text-decoration:underline;">🗑️ Borrar chat</button>
      </div>
    </div>

    <!-- ── LECCIÓN ── -->
    <div id="panel-leccion" class="ingles-panel" style="display:none;">
      <div style="margin-bottom:10px;color:#00cc66;font-size:13px;">Clases estructuradas con ejercicios y evaluación. 🎓</div>
      <div id="inglesLeccionChat" style="background:#000d06;border:1px solid #00441a;border-radius:12px;padding:14px;height:300px;overflow-y:auto;margin-bottom:12px;display:flex;flex-direction:column;gap:10px;"></div>
      <div style="display:flex;gap:8px;align-items:flex-end;">
        <textarea id="inglesLeccionInput" placeholder="Respondé los ejercicios..." style="flex:1;background:#001a0d;color:#e6ffe6;border:1px solid #00cc6644;border-radius:10px;padding:10px 12px;font-size:14px;resize:none;height:52px;outline:none;font-family:'Segoe UI',sans-serif;" onkeydown="if(event.key==='Enter'&&!event.shiftKey){event.preventDefault();inglesEnviarLeccion();}"></textarea>
        <button onclick="inglesEnviarLeccion()" style="background:linear-gradient(135deg,#003d1a,#006633);color:#00ff88;border:1px solid #00cc66;border-radius:10px;padding:12px 18px;cursor:pointer;font-size:18px;">➤</button>
      </div>
      <div style="margin-top:10px;display:flex;gap:8px;flex-wrap:wrap;">
        <button onclick="inglesNuevaLeccion()" style="background:#002a1a;border:1px solid #00cc6644;color:#00ff88;border-radius:8px;padding:8px 16px;cursor:pointer;font-size:13px;">🎓 Nueva lección</button>
        <button onclick="inglesCompletarLeccion()" style="background:#003d1a;border:1px solid #00cc66;color:#00ff88;border-radius:8px;padding:8px 16px;cursor:pointer;font-size:13px;">✅ Completar (+50 pts)</button>
      </div>
    </div>

    <!-- ── ESCENARIO ── -->
    <div id="panel-escenario" class="ingles-panel" style="display:none;">
      <div style="margin-bottom:10px;color:#00cc66;font-size:13px;">Simulá situaciones reales en inglés. El profesor interpreta el rol y corrige. 🌍</div>
      <div id="escenarioSelector" style="display:grid;grid-template-columns:repeat(auto-fill,minmax(150px,1fr));gap:10px;margin-bottom:14px;"></div>
      <div id="escenarioChat" style="display:none;">
        <div id="escenarioChatNombre" style="color:#00ff88;font-weight:700;margin-bottom:8px;font-size:14px;"></div>
        <div id="inglesEscenarioMessages" style="background:#000d06;border:1px solid #00441a;border-radius:12px;padding:14px;height:280px;overflow-y:auto;margin-bottom:12px;display:flex;flex-direction:column;gap:10px;"></div>
        <div style="display:flex;gap:8px;align-items:flex-end;">
          <textarea id="inglesEscenarioInput" placeholder="Respond in English..." style="flex:1;background:#001a0d;color:#e6ffe6;border:1px solid #00cc6644;border-radius:10px;padding:10px 12px;font-size:14px;resize:none;height:52px;outline:none;font-family:'Segoe UI',sans-serif;" onkeydown="if(event.key==='Enter'&&!event.shiftKey){event.preventDefault();inglesEnviarEscenario();}"></textarea>
          <button onclick="inglesEnviarEscenario()" style="background:linear-gradient(135deg,#003d1a,#006633);color:#00ff88;border:1px solid #00cc66;border-radius:10px;padding:12px 18px;cursor:pointer;font-size:18px;">➤</button>
        </div>
        <button onclick="volverEscenarios()" style="margin-top:10px;background:transparent;border:none;color:#00cc66;font-size:13px;cursor:pointer;text-decoration:underline;">← Elegir otro escenario</button>
      </div>
    </div>

    <!-- ── PROGRESO ── -->
    <div id="panel-progreso" class="ingles-panel" style="display:none;">
      <div id="inglesStats" style="background:#000d06;border:1px solid #00441a;border-radius:14px;padding:18px;margin-bottom:14px;font-size:14px;line-height:1.9;white-space:pre-wrap;color:#e6ffe6;">Cargando...</div>
      <div style="background:#001a0d;border:1px solid #00cc6622;border-radius:12px;padding:14px;margin-bottom:12px;">
        <div style="color:#00ff88;font-weight:700;margin-bottom:10px;">🎯 Tu nivel CEFR</div>
        <div id="nivelesGrid" style="display:grid;grid-template-columns:repeat(3,1fr);gap:8px;"></div>
      </div>
      <div style="background:#001a0d;border:1px solid #00cc6622;border-radius:12px;padding:14px;">
        <div style="color:#00ff88;font-weight:700;margin-bottom:8px;">🏆 Logros</div>
        <div id="inglesLogros" style="color:#00aa55;font-size:13px;line-height:1.9;">Cargando...</div>
      </div>
    </div>

  </div>
</div>

<style>
.ingles-tab{background:#001a0d;border:1px solid #00cc6644;color:#00aa55;border-radius:8px;padding:7px 16px;cursor:pointer;font-size:13px;font-weight:600;transition:all .2s;}
.ingles-tab.active,.ingles-tab:hover{background:#003d1a;border-color:#00cc66;color:#00ff88;}
.ingles-burbuja-user{align-self:flex-end;background:#003d1a;border:1px solid #00cc6633;color:#e6ffe6;border-radius:14px 14px 4px 14px;padding:9px 14px;max-width:85%;font-size:13px;line-height:1.5;word-break:break-word;}
.ingles-burbuja-ai{align-self:flex-start;background:#001a0d;border:1px solid #00441a;color:#e6ffe6;border-radius:4px 14px 14px 14px;padding:9px 14px;max-width:90%;font-size:13px;line-height:1.7;word-break:break-word;white-space:pre-wrap;}
</style>

<script>
// ══════════════════════════════════════════════
//  📚 PROFESOR DE INGLÉS — JS
// ══════════════════════════════════════════════
let inglesTabActual = "conversacion";
let inglesEscenarioActual = "";
let inglesPerfil = null;

function abrirProfesorIngles() {
  document.getElementById("modalIngles").style.display = "flex";
  inglesCargarPerfil();
}
function cerrarIngles() {
  document.getElementById("modalIngles").style.display = "none";
}

function inglesSetTab(tab) {
  inglesTabActual = tab;
  document.querySelectorAll(".ingles-tab").forEach(b => b.classList.remove("active"));
  document.querySelectorAll(".ingles-panel").forEach(p => p.style.display = "none");
  document.getElementById("tab-" + tab).classList.add("active");
  document.getElementById("panel-" + tab).style.display = "block";
  if (tab === "progreso") inglesCargarProgreso();
  if (tab === "escenario") inglesRenderEscenarios();
}

async function inglesCargarPerfil() {
  const r = await fetch("/ingles/perfil");
  const d = await r.json();
  inglesPerfil = d;
  document.getElementById("inglesNivelBadge").textContent =
    d.perfil.nivel + " — " + (d.niveles[d.perfil.nivel]?.nombre || "");
}

// ── Conversación ───────────────────────────────
async function inglesEnviar() {
  const input = document.getElementById("inglesInput");
  const msg = input.value.trim();
  if (!msg) return;
  input.value = "";
  inglesMsg("inglesChat", msg, "user");
  const load = inglesMsg("inglesChat", "✍️ Thinking...", "ai");
  const r = await fetch("/ingles/chat", {
    method: "POST", headers: {"Content-Type":"application/json"},
    body: JSON.stringify({mensaje: msg, sub_modo: "conversacion"})
  });
  const d = await r.json();
  load.remove();
  inglesMsg("inglesChat", d.texto, "ai");
}

function inglesHablar() {
  const Rec = window.SpeechRecognition || window.webkitSpeechRecognition;
  if (!Rec) { alert("Tu navegador no soporta reconocimiento de voz."); return; }
  const rec = new Rec();
  rec.lang = "en-US"; rec.interimResults = false;
  const btn = document.getElementById("inglesVozBtn");
  btn.textContent = "🔴"; rec.start();
  rec.onresult = e => { document.getElementById("inglesInput").value = e.results[0][0].transcript; btn.textContent = "🎤"; inglesEnviar(); };
  rec.onerror = rec.onend = () => btn.textContent = "🎤";
}

// ── Lección ────────────────────────────────────
function inglesNuevaLeccion() {
  document.getElementById("inglesLeccionChat").innerHTML = "";
  inglesMsg("inglesLeccionChat", "Hello! I'm your English teacher. Let's start a new lesson. Ready? Escribí 'start' para comenzar. 🎓", "ai");
}
async function inglesEnviarLeccion() {
  const input = document.getElementById("inglesLeccionInput");
  const msg = input.value.trim(); if (!msg) return; input.value = "";
  inglesMsg("inglesLeccionChat", msg, "user");
  const load = inglesMsg("inglesLeccionChat", "✍️ Preparando lección...", "ai");
  const r = await fetch("/ingles/chat", { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({mensaje:msg, sub_modo:"leccion"}) });
  const d = await r.json(); load.remove();
  inglesMsg("inglesLeccionChat", d.texto, "ai");
}
async function inglesCompletarLeccion() {
  const r = await fetch("/ingles/leccion_completada", {method:"POST"});
  const d = await r.json();
  if (d.ok) { inglesMsg("inglesLeccionChat", "🏆 ¡Lección completada! +50 pts y 10 palabras.\n\n" + d.resumen, "ai"); inglesCargarPerfil(); }
}

// ── Escenarios ─────────────────────────────────
function inglesRenderEscenarios() {
  if (!inglesPerfil) { inglesCargarPerfil().then(inglesRenderEscenarios); return; }
  const grid = document.getElementById("escenarioSelector");
  grid.innerHTML = "";
  Object.entries(inglesPerfil.escenarios || {}).forEach(([key, esc]) => {
    const btn = document.createElement("button");
    btn.style.cssText = "background:#001a0d;border:1px solid #00cc6633;border-radius:12px;padding:14px 10px;cursor:pointer;color:#00ff88;font-size:13px;font-weight:600;text-align:center;transition:all .2s;width:100%;";
    btn.innerHTML = `<div style="font-size:26px;margin-bottom:6px;">${esc.emoji}</div>${esc.nombre}`;
    btn.onmouseover = () => btn.style.borderColor = "#00cc66";
    btn.onmouseout  = () => btn.style.borderColor = "#00cc6633";
    btn.onclick = () => inglesIniciarEscenario(key, esc);
    grid.appendChild(btn);
  });
}
async function inglesIniciarEscenario(key, esc) {
  inglesEscenarioActual = key;
  document.getElementById("escenarioSelector").style.display = "none";
  document.getElementById("escenarioChat").style.display = "block";
  document.getElementById("escenarioChatNombre").textContent = esc.emoji + " " + esc.nombre;
  document.getElementById("inglesEscenarioMessages").innerHTML = "";
  const r = await fetch("/ingles/chat", { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({mensaje:"__inicio__", sub_modo:"escenario", escenario:key}) });
  const d = await r.json();
  inglesMsg("inglesEscenarioMessages", d.texto, "ai");
}
async function inglesEnviarEscenario() {
  const input = document.getElementById("inglesEscenarioInput");
  const msg = input.value.trim(); if (!msg) return; input.value = "";
  inglesMsg("inglesEscenarioMessages", msg, "user");
  const load = inglesMsg("inglesEscenarioMessages", "✍️ Responding...", "ai");
  const r = await fetch("/ingles/chat", { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({mensaje:msg, sub_modo:"escenario", escenario:inglesEscenarioActual}) });
  const d = await r.json(); load.remove();
  inglesMsg("inglesEscenarioMessages", d.texto, "ai");
}
function volverEscenarios() {
  inglesEscenarioActual = "";
  document.getElementById("escenarioSelector").style.display = "grid";
  document.getElementById("escenarioChat").style.display = "none";
}

// ── Progreso ───────────────────────────────────
async function inglesCargarProgreso() {
  const r = await fetch("/ingles/perfil");
  const d = await r.json(); inglesPerfil = d;
  document.getElementById("inglesStats").textContent = d.resumen;
  const grid = document.getElementById("nivelesGrid"); grid.innerHTML = "";
  Object.entries(d.niveles).forEach(([cod, info]) => {
    const btn = document.createElement("button");
    const activo = cod === d.perfil.nivel;
    btn.style.cssText = `background:${activo?"#003d1a":"#001a0d"};border:1px solid ${activo?"#00cc66":"#00441a"};border-radius:8px;padding:10px 6px;cursor:pointer;color:${activo?"#00ff88":"#00aa55"};font-size:12px;font-weight:${activo?"800":"600"};text-align:center;`;
    btn.innerHTML = `${info.emoji}<br><b>${cod}</b><br><span style="font-size:10px;opacity:.8">${info.nombre}</span>`;
    btn.onclick = () => inglesSetNivel(cod);
    grid.appendChild(btn);
  });
  const logros = d.perfil.logros || [];
  document.getElementById("inglesLogros").innerHTML = logros.length
    ? logros.map(l => `<div>✅ ${l}</div>`).join("")
    : "Todavía no desbloqueaste logros. ¡Empezá a estudiar! 💪";
}
async function inglesSetNivel(nivel) {
  const r = await fetch("/ingles/cambiar_nivel", { method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({nivel}) });
  const d = await r.json();
  if (d.ok) { document.getElementById("inglesNivelBadge").textContent = d.nivel + " — " + d.nombre; inglesCargarProgreso(); }
}

// ── Utilidad ───────────────────────────────────
function inglesMsg(chatId, texto, tipo) {
  const chat = document.getElementById(chatId);
  const div = document.createElement("div");
  div.className = tipo === "user" ? "ingles-burbuja-user" : "ingles-burbuja-ai";
  div.textContent = texto;
  chat.appendChild(div);
  chat.scrollTop = chat.scrollHeight;
  return div;
}
</script>
   
</body>
</html>

"""

# ---------------- RUTAS ----------------
import mercadopago

sdk = mercadopago.SDK(os.getenv("MP_ACCESS_TOKEN"))

@app.route("/auth/register", methods=["POST"])
def register():
    data = request.get_json()

    email = data.get("email", "").lower().strip()
    password = data.get("password", "")

    ok, msg = registrar_usuario(email, password)
    if not ok:
        return jsonify({"ok": False, "msg": msg})

    # login automático
    session["user_email"] = email
    return jsonify({"ok": True})

@app.route("/auth/login", methods=["POST"])
def login():
    data = request.get_json()

    email = data.get("email", "").lower().strip()
    password = data.get("password", "")

    if not autenticar_usuario(email, password):
        return jsonify({"ok": False, "msg": "Credenciales incorrectas"})

    session["user_email"] = email
    return jsonify({"ok": True})

@app.route("/auth/logout")
def logout():
    session.pop("user_email", None)
    return redirect("/")

@app.route("/premium")
def premium():

    # 1️⃣ Verificar login
    usuario = session.get("user_email")
    if not usuario:
        return jsonify({"error": "No logueado"}), 401

    # 2️⃣ Determinar tipo de plan
    tipo = request.args.get("tipo", "mensual")

    if tipo == "anual":
        titulo = "Foschi IA Premium Anual (12 meses PAGA 10)"
        precio = 100000   # 🔥 10 meses pagos
    else:
        titulo = "Foschi IA Premium Mensual"
        precio = 10000

    # 3️⃣ Crear preferencia MercadoPago
    pref = {
        "items": [{
            "title": titulo,
            "quantity": 1,
            "unit_price": precio
        }],
        "external_reference": usuario,
        "notification_url": "https://foschi-ia.onrender.com/webhook/mp"
    }

    preference = sdk.preference().create(pref)
    return jsonify(preference["response"])

@app.route("/webhook/mp", methods=["POST"])
def webhook_mp():
    data = request.json

    if not data or "data" not in data:
        return "ok"

    if data.get("type") != "payment":
        return "ok"

    payment_id = data["data"].get("id")
    if not payment_id:
        return "ok"

    payment = sdk.payment().get(payment_id)
    info = payment.get("response", {})

    # 🔐 VALIDAR QUE EL PAGO SEA TUYO
    if info.get("merchant_account_id") != os.getenv("MP_MERCHANT_ID"):
        return "ok"

    if info.get("status") != "approved":
        return "ok"

    usuario = info.get("external_reference")
    if not usuario:
        return "ok"

    from pagos import pago_ya_registrado, registrar_pago
    if pago_ya_registrado(str(payment_id)):
        return "ok"

    items = info.get("additional_info", {}).get("items", [])
    titulo = items[0]["title"] if items else ""
    plan = "anual" if "12" in titulo or "Anual" in titulo else "mensual"

    monto = info.get("transaction_amount", 0)

    activar_premium(usuario, plan)

    registrar_pago(
        usuario=usuario,
        monto=monto,
        plan=plan,
        payment_id=str(payment_id)
    )

    return "ok"

@app.route("/")
def index():
    if "usuario_id" not in session:
        session["usuario_id"] = str(uuid.uuid4())

    usuario = session.get("user_email") or session["usuario_id"]

    premium = usuario_premium(usuario)

    is_super = es_superusuario(usuario)
    rol = rol_superusuario(usuario)
    nivel = nivel_superusuario(usuario)

    return render_template_string(
        HTML_TEMPLATE,
        APP_NAME=APP_NAME,
        usuario_id=usuario,
        premium=premium,
        is_super=is_super,
        rol=rol,
        nivel=nivel
    )

@app.route("/preguntar", methods=["POST"])
def preguntar():

    data = request.get_json()

    mensaje = data.get("mensaje", "")

    doc_id = data.get("doc_id")

    preguntar_doc = data.get("preguntar_doc", False)
    
    # 1️⃣ Asegurar UUID en sesión (NO se borra nunca)
    if "usuario_id" not in session:
        session["usuario_id"] = str(uuid.uuid4())

    # 2️⃣ Identidad activa:
    #    email si está logueado
    #    UUID si no
    usuario = session.get("user_email") or session["usuario_id"]

    lat = data.get("lat")
    lon = data.get("lon")
    tz  = data.get("timeZone") or data.get("time_zone") or None

     # ============================
    # 📄 PREGUNTAS SOBRE DOCUMENTO
    # ============================

    if preguntar_doc and doc_id:

        txt_path = os.path.join(TEMP_DIR, f"{doc_id}.txt")

        if os.path.exists(txt_path):

            with open(txt_path, "r", encoding="utf-8") as f:
                contenido_doc = f.read()

            try:

                client = OpenAI(api_key=OPENAI_API_KEY)

                prompt = f"""
Sos Foschi IA.

Respondé usando SOLAMENTE el contenido del documento.

DOCUMENTO:
{contenido_doc[:12000]}

PREGUNTA:
{mensaje}
"""

                resp = client.chat.completions.create(
                    model="gpt-4-turbo",
                    messages=[
                        {
                            "role": "user",
                            "content": prompt
                        }
                    ],
                    temperature=0.2,
                    max_tokens=700
                )

                texto = resp.choices[0].message.content.strip()

                return jsonify({
                    "texto": texto,
                    "imagenes": [],
                    "borrar_historial": False
                })

            except Exception as e:

                return jsonify({
                    "texto": f"Error analizando documento: {e}",
                    "imagenes": [],
                    "borrar_historial": False
                })

    # 3️⃣ Generar respuesta con identidad correcta
    respuesta = generar_respuesta(
        mensaje,
        usuario,
        lat=lat,
        lon=lon,
        tz=tz
    )

    # 4️⃣ Guardar historial con la MISMA identidad
    texto_para_hist = (
        respuesta["texto"]
        if isinstance(respuesta, dict) and "texto" in respuesta
        else str(respuesta)
    )

    guardar_en_historial(usuario, mensaje, texto_para_hist)

    return jsonify(respuesta)

@app.route("/historial/<usuario_id>")
def historial(usuario_id):
    return jsonify(cargar_historial(usuario_id))

# ──────────────────────────────────────────────
# 🇬🇧  PROFESOR DE INGLÉS — RUTAS
# ──────────────────────────────────────────────

@app.route("/ingles/perfil")
def ingles_ver_perfil():
    if "usuario_id" not in session:
        session["usuario_id"] = str(uuid.uuid4())
    usuario = session.get("user_email") or session["usuario_id"]
    perfil = ingles_perfil(usuario)
    return jsonify({
        "perfil": perfil,
        "resumen": ingles_resumen(usuario),
        "niveles": NIVELES,
        "escenarios": {k: {"emoji": v["emoji"], "nombre": v["nombre"]} for k, v in ESCENARIOS.items()},
    })

@app.route("/ingles/chat", methods=["POST"])
def ingles_chat():
    data = request.get_json()
    mensaje   = (data.get("mensaje") or "").strip()
    sub_modo  = data.get("sub_modo", "conversacion")
    escenario = data.get("escenario", "")

    if "usuario_id" not in session:
        session["usuario_id"] = str(uuid.uuid4())
    usuario = session.get("user_email") or session["usuario_id"]

    perfil = ingles_perfil(usuario)
    nivel  = perfil.get("nivel", "A1")
    ingles_racha(usuario)

    if sub_modo == "leccion":
        system = prompt_leccion(nivel)
    elif sub_modo == "escenario" and escenario:
        system = prompt_escenario(escenario, nivel)
    else:
        system = prompt_conversacion_basica(nivel, perfil)

    try:
        client_ai = OpenAI(api_key=OPENAI_API_KEY)
        resp = client_ai.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": system},
                {"role": "user",   "content": mensaje},
            ],
            temperature=0.65,
            max_tokens=700,
        )
        texto = resp.choices[0].message.content.strip()
        if "✅ Correcto:" in texto or "✅ Correction:" in texto:
            ingles_error(usuario, mensaje[:80])
            ingles_puntos(usuario, 5)
        else:
            ingles_puntos(usuario, 10)
        if "📖" in texto or "nueva palabra" in texto.lower():
            ingles_puntos(usuario, 0, palabras=1)
    except Exception as e:
        texto = f"Error al conectar con la IA: {e}"

    return jsonify({"texto": texto})

@app.route("/ingles/leccion_completada", methods=["POST"])
def ingles_leccion_completada():
    if "usuario_id" not in session:
        session["usuario_id"] = str(uuid.uuid4())
    usuario = session.get("user_email") or session["usuario_id"]
    ingles_completar(usuario)
    ingles_puntos(usuario, 50, palabras=10)
    return jsonify({"ok": True, "resumen": ingles_resumen(usuario)})

@app.route("/ingles/cambiar_nivel", methods=["POST"])
def ingles_cambiar_nivel():
    data  = request.get_json()
    nivel = data.get("nivel", "A1")
    if nivel not in NIVELES:
        return jsonify({"ok": False, "msg": "Nivel inválido"})
    if "usuario_id" not in session:
        session["usuario_id"] = str(uuid.uuid4())
    usuario = session.get("user_email") or session["usuario_id"]
    perfil = ingles_perfil(usuario)
    perfil["nivel"] = nivel
    ingles_guardar(usuario, perfil)
    return jsonify({"ok": True, "nivel": nivel, "nombre": NIVELES[nivel]["nombre"]})


def tts():
    texto = request.args.get("texto","")
    try:
        tts_obj = gTTS(text=texto, lang="es", slow=False, tld="com.mx")
        archivo = io.BytesIO()
        tts_obj.write_to_fp(archivo)
        archivo.seek(0)
        return send_file(archivo, mimetype="audio/mpeg")
    except Exception as e:
        return f"Error TTS: {e}", 500

@app.route("/clima")
def clima():
    lat = request.args.get("lat")
    lon = request.args.get("lon")
    ciudad = request.args.get("ciudad")
    return obtener_clima(ciudad=ciudad, lat=lat, lon=lon)

@app.route('/favicon.ico')
def favicon():
    ico = os.path.join(STATIC_DIR, 'favicon.ico')
    if os.path.exists(ico):
        return send_file(ico)
    return "", 204

@app.route("/avisos", methods=["POST"])
def avisos():
    usuario = request.json.get("usuario_id", "anon")
    lista = load_recordatorios()
    ahora = datetime.now(TZ)
    vencidos = []
    restantes = []
    for r in lista:
        when = TZ.localize(datetime.strptime(r["cuando"], "%Y-%m-%d %H:%M:%S"))
        if when <= ahora and r["usuario"] == usuario:
            vencidos.append(r)
        else:
            restantes.append(r)
    save_recordatorios(restantes)
    return jsonify(vencidos)

@app.route("/admin/pagos")
def admin_pagos():
    import json, os

    # 🔐 clave simple (después se mejora)
    if request.args.get("key") != "foschi_admin_2026":
        return "Acceso denegado", 403

    archivo = "data/pagos.json"
    if not os.path.exists(archivo):
        return "<h2>No hay pagos todavía</h2>"

    pagos = json.load(open(archivo))

    html = """
    <h2>💎 Pagos Foschi IA</h2>
    <table border="1" cellpadding="8">
      <tr>
        <th>Usuario</th>
        <th>Plan</th>
        <th>Fecha</th>
        <th>Payment ID</th>
        <th>Status</th>
      </tr>
    """

    for u, p in pagos.items():
        html += f"""
        <tr>
          <td>{u}</td>
          <td>{p['plan']}</td>
          <td>{p['fecha']}</td>
          <td>{p['payment_id']}</td>
          <td>{p['status']}</td>
        </tr>
        """

    html += "</table>"
    return html

# ---------------- AUDIO A WORD DOCX ----------------
from werkzeug.utils import secure_filename

@app.route("/upload_audio", methods=["POST"])
def upload_audio():
    if "audio" not in request.files:
        return "No se envió archivo", 400
    
    file = request.files["audio"]
    usuario_id = request.form.get("usuario_id", "anon")

    # Guardar archivo temporal
    filename = secure_filename(file.filename)
    if not filename:
        return "Nombre de archivo inválido", 400
    temp_path = os.path.join("temp", f"{uuid.uuid4()}_{filename}")
    os.makedirs("temp", exist_ok=True)
    file.save(temp_path)

    docx_path = None  # PREVENIR ERROR EN finally

    try:
        # ---- TRANSCRIPCIÓN OPENAI ----
        with open(temp_path, "rb") as f:
            transcript = client.audio.transcriptions.create(
                model="gpt-4o-transcribe",
                file=f,
            )

        texto_transcrito = transcript.text if hasattr(transcript, "text") else str(transcript)

        # ---- CREAR DOCX ----
        nombre_docx = filename.rsplit(".", 1)[0] + ".docx"
        docx_path = os.path.join("temp", nombre_docx)

        doc = DocxDocument()
        doc.add_heading("Transcripción de audio", level=1)
        doc.add_paragraph(texto_transcrito)
        doc.add_page_break()
        doc.save(docx_path)

        # --- programar limpieza después de responder ---
        @after_this_request
        def _cleanup(response):
            try:
                if os.path.exists(docx_path):
                    os.remove(docx_path)
            except Exception:
                pass
            try:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
            except Exception:
                pass
            return response

        # ---- ENVIAR DOCX ----
        return send_file(
            docx_path,
            as_attachment=True,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            download_name=nombre_docx
        )

    except Exception as e:
        # intentar limpiar si algo quedó
        try:
            if docx_path and os.path.exists(docx_path):
                os.remove(docx_path)
        except:
            pass
        try:
            if os.path.exists(temp_path):
                os.remove(temp_path)
        except:
            pass
        return f"Error en transcripción: {str(e)}", 500

# ---------------- NUEVOS ENDPOINTS: subir documento (extraer texto) y resumir (crear .docx) ----------------
def extract_text_from_pdf(path):
    text = ""
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                try:
                    p = page.extract_text()
                    if p:
                        text += p + "\n"
                except:
                    continue
    except Exception as e:
        print("Error leyendo PDF:", e)
    return text

def extract_text_from_docx(path):
    text = ""
    try:
        doc = docx_reader.Document(path)
        for p in doc.paragraphs:
            if p.text:
                text += p.text + "\n"
    except Exception as e:
        print("Error leyendo DOCX:", e)
    return text

@app.route("/upload_doc", methods=["POST"])
def upload_doc():
    """Recibe PDF o DOCX, extrae texto y guarda temporalmente. Devuelve doc_id que luego se usa para pedir resumen."""
    if "archivo" not in request.files:
        return "No se envió archivo", 400
    file = request.files["archivo"]
    usuario_id = request.form.get("usuario_id", "anon")
    filename = secure_filename(file.filename)
    if filename == "":
        return "Archivo sin nombre", 400
    ext = filename.rsplit(".",1)[-1].lower()
    if ext not in ["pdf", "docx"]:
        return "Formato no permitido. Solo PDF o DOCX.", 400

    doc_id = str(uuid.uuid4())
    saved_name = f"{doc_id}_{filename}"
    temp_path = os.path.join(TEMP_DIR, saved_name)
    try:
        file.save(temp_path)
    except Exception as e:
        return f"Error guardando archivo temporal: {e}", 500

    # extraer texto
    if ext == "pdf":
        text = extract_text_from_pdf(temp_path)
    else:
        text = extract_text_from_docx(temp_path)

    if not text or len(text.strip()) == 0:
        # limpiar archivo
        try:
            os.remove(temp_path)
        except:
            pass
        return "No pude extraer texto del documento.", 400

    # guardar texto en fichero temporal .txt
    txt_path = os.path.join(TEMP_DIR, f"{doc_id}.txt")
    try:
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(text)
    except Exception as e:
        # limpiar archivo original
        try:
            os.remove(temp_path)
        except:
            pass
        return f"Error guardando texto temporal: {e}", 500

    # devolvemos doc_id y un snippet para mostrar
    snippet = text[:800].replace("\n"," ") + ("..." if len(text)>800 else "")
    return jsonify({"doc_id": doc_id, "name": filename, "snippet": snippet})

@app.route("/resumir_doc", methods=["POST"])
def resumir_doc():

    data = request.get_json()

    doc_id = data.get("doc_id")

    modo = data.get("modo", "normal")

    txt_path = os.path.join(TEMP_DIR, f"{doc_id}.txt")

    if not os.path.exists(txt_path):
        return "Documento no encontrado", 404

    with open(txt_path, "r", encoding="utf-8") as f:
        texto = f.read()

    # ============================
    # TIPOS DE RESUMEN
    # ============================

    if modo == "breve":

        instrucciones = (
            "Hacé un resumen breve y directo "
            "del siguiente documento."
        )

    elif modo == "profundo":

        instrucciones = (
            "Hacé un resumen MUY detallado "
            "del siguiente documento. "
            "Separá por temas y explicá bien."
        )

    else:

        instrucciones = (
            "Resumí el siguiente texto "
            "de forma clara, ordenada y completa. "
            "Usá títulos y viñetas si hace falta."
        )

    # ============================
    # GENERAR RESUMEN IA
    # ============================

    try:

        client = OpenAI(api_key=OPENAI_API_KEY)

        prompt = f"""
{instrucciones}

TEXTO:
{texto[:15000]}
"""

        resp = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.4,
            max_tokens=1500
        )

        resumen = resp.choices[0].message.content.strip()

    except Exception as e:

        return f"Error generando resumen: {e}", 500

    # ============================
    # CREAR WORD
    # ============================

    nombre_doc = f"resumen_{doc_id}.docx"

    ruta_doc = os.path.join(TEMP_DIR, nombre_doc)

    doc = DocxDocument()

    doc.add_heading(
        "Resumen generado por Foschi IA",
        level=1
    )

    doc.add_paragraph(resumen)

    doc.save(ruta_doc)

    @after_this_request
    def cleanup(response):

        try:

            if os.path.exists(ruta_doc):
                os.remove(ruta_doc)

        except:
            pass

        return response

    return send_file(
        ruta_doc,
        as_attachment=True,
        download_name="resumen_foschi.docx"
    )

# ---------------- GENERADOR DE PRESENTACIONES (PPTX) ----------------

# Generar la presentación (texto IA + imágenes IA) puede tardar más de lo que
# permite el timeout del servidor/proxy si se hace todo en una sola request
# (eso provoca errores 500/502). Por eso el proceso se hace en un hilo en
# segundo plano y el front-end consulta el estado por separado.
PRESENTACIONES_JOBS = {}
PRESENTACIONES_LOCK = threading.Lock()


def generar_estructura_presentacion(contenido_base, tema, num_slides=8):
    """
    Usa OpenAI para generar la estructura de una presentación en formato JSON.
    Devuelve un dict con 'titulo_presentacion', 'subtitulo' y 'diapositivas'
    (lista de {'titulo','bullets','notas','imagen_prompt'}), o None si falla.
    """
    try:
        cliente = OpenAI(api_key=OPENAI_API_KEY)

        if contenido_base and contenido_base.strip():
            fuente = (
                "Basate en el siguiente contenido/documento para armar la presentación:\n\n"
                f"{contenido_base[:12000]}\n\n"
                f"Enfoque adicional pedido por el usuario (si aplica): {tema or 'usar el contenido tal cual, resumido'}"
            )
        else:
            fuente = f"Tema de la presentación: {tema}"

        prompt = f"""
Sos un diseñador experto de presentaciones corporativas. Generá la estructura de una
presentación de PowerPoint de {num_slides} diapositivas (sin contar la portada), en español.

{fuente}

Respondé EXCLUSIVAMENTE con un JSON válido, sin texto adicional, sin explicaciones y sin
bloques de markdown (sin ```), con esta forma exacta:

{{
  "titulo_presentacion": "Título principal de la presentación",
  "subtitulo": "Subtítulo o frase introductoria corta",
  "diapositivas": [
    {{
      "titulo": "Título de la diapositiva",
      "bullets": ["punto 1", "punto 2", "punto 3"],
      "notas": "Notas para el orador (1-2 frases)",
      "imagen_prompt": "Descripción corta EN INGLÉS para generar una imagen ilustrativa con IA, sin texto dentro de la imagen"
    }}
  ]
}}

Reglas:
- Cada diapositiva debe tener entre 3 y 5 bullets, breves, claros y sin repetir el título.
- Generá exactamente {num_slides} diapositivas dentro de "diapositivas".
- No incluyas comentarios ni texto fuera del JSON.
"""

        resp = cliente.chat.completions.create(
            model="gpt-4-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.6,
            max_tokens=3500
        )

        texto = resp.choices[0].message.content.strip()

        # Limpiar posibles bloques de markdown ```json ... ```
        texto = re.sub(r"^```(json)?", "", texto.strip())
        texto = re.sub(r"```$", "", texto.strip())
        texto = texto.strip()

        # Por si OpenAI agrega texto extra antes/después del JSON
        inicio = texto.find("{")
        fin = texto.rfind("}")
        if inicio != -1 and fin != -1 and fin > inicio:
            texto = texto[inicio:fin + 1]

        estructura = json.loads(texto)

        if not isinstance(estructura.get("diapositivas"), list) or not estructura["diapositivas"]:
            return None

        return estructura

    except Exception as e:
        print("Error generando estructura de presentación:", e)
        traceback.print_exc()
        return None


def generar_imagen_presentacion_bytes(prompt_imagen):
    """Genera una imagen con IA para una diapositiva. Devuelve bytes PNG o None si falla."""
    try:
        cliente = OpenAI(api_key=OPENAI_API_KEY)
        resultado = cliente.images.generate(
            model="gpt-image-1",
            prompt=prompt_imagen or "ilustración abstracta minimalista, colores azules",
            size="1024x1024",
            quality="medium"
        )
        b64 = resultado.data[0].b64_json
        return base64.b64decode(b64)
    except Exception as e:
        print("Error generando imagen para presentación:", e)
        traceback.print_exc()
        return None


def _agregar_fondo(slide, prs, color):
    """Agrega un rectángulo de fondo de color sólido a toda la diapositiva."""
    fondo = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = color
    fondo.line.fill.background()
    fondo.shadow.inherit = False
    # Enviar el fondo detrás de los demás elementos
    spTree = slide.shapes._spTree
    spTree.remove(fondo._element)
    spTree.insert(2, fondo._element)
    return fondo


def construir_pptx(estructura, incluir_imagenes=True, video_paths=None):
    """
    Construye un archivo .pptx a partir de la estructura generada por IA.
    video_paths: lista de rutas a archivos de video cortos (opcional) para
    insertar en las primeras diapositivas en lugar de imágenes generadas.
    Devuelve la ruta del archivo .pptx generado (en TEMP_DIR).
    """
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    COLOR_FONDO = RGBColor(0x00, 0x14, 0x24)
    COLOR_TITULO = RGBColor(0x00, 0xEA, 0xFF)
    COLOR_TEXTO = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_ACENTO = RGBColor(0x33, 0xAA, 0xCC)

    blank_layout = prs.slide_layouts[6]

    # ---- Diapositiva de portada ----
    slide = prs.slides.add_slide(blank_layout)
    _agregar_fondo(slide, prs, COLOR_FONDO)

    titulo_box = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(2.6), PptxInches(11.3), PptxInches(1.8)
    )
    tf = titulo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = estructura.get("titulo_presentacion") or "Presentación"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    p.alignment = PP_ALIGN.CENTER

    sub = estructura.get("subtitulo") or ""
    if sub:
        sub_box = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(4.3), PptxInches(11.3), PptxInches(1)
        )
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = sub
        p2.font.size = Pt(22)
        p2.font.color.rgb = COLOR_TEXTO
        p2.alignment = PP_ALIGN.CENTER

    pie = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(6.7), PptxInches(11.3), PptxInches(0.5)
    )
    pf = pie.text_frame.paragraphs[0]
    pf.text = f"Generado por {APP_NAME}"
    pf.font.size = Pt(12)
    pf.font.color.rgb = COLOR_ACENTO
    pf.alignment = PP_ALIGN.CENTER

    # ---- Diapositivas de contenido ----
    diapositivas = estructura.get("diapositivas", [])
    video_paths = video_paths or []
    video_idx = 0

    # Pre-generar TODAS las imágenes en paralelo (si corresponde) para no
    # acumular tiempos de espera secuenciales y evitar timeouts del servidor.
    imagenes_por_slide = [None] * len(diapositivas)
    if incluir_imagenes:
        indices_a_generar = []
        prompts_a_generar = []
        for idx, dia in enumerate(diapositivas):
            if idx < len(video_paths):
                continue  # esa diapositiva usará video, no imagen
            img_prompt = dia.get("imagen_prompt") or dia.get("titulo") or estructura.get("titulo_presentacion", "")
            indices_a_generar.append(idx)
            prompts_a_generar.append(img_prompt)

        if prompts_a_generar:
            try:
                with ThreadPoolExecutor(max_workers=min(5, len(prompts_a_generar))) as executor:
                    resultados = list(executor.map(generar_imagen_presentacion_bytes, prompts_a_generar))
                for idx, img_bytes in zip(indices_a_generar, resultados):
                    imagenes_por_slide[idx] = img_bytes
            except Exception:
                print("Error generando imágenes en paralelo:")
                traceback.print_exc()

    for idx, dia in enumerate(diapositivas):
        slide = prs.slides.add_slide(blank_layout)
        _agregar_fondo(slide, prs, COLOR_FONDO)

        hay_video = video_idx < len(video_paths)
        hay_media = hay_video or incluir_imagenes
        ancho_texto = PptxInches(7.0) if hay_media else PptxInches(11.7)

        # Título
        titulo_box = slide.shapes.add_textbox(
            PptxInches(0.6), PptxInches(0.4), PptxInches(12.1), PptxInches(1)
        )
        tf = titulo_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = dia.get("titulo") or f"Diapositiva {idx + 1}"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITULO

        # Línea decorativa bajo el título
        linea = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, PptxInches(0.6), PptxInches(1.35), PptxInches(3.5), Pt(3)
        )
        linea.fill.solid()
        linea.fill.fore_color.rgb = COLOR_TITULO
        linea.line.fill.background()
        linea.shadow.inherit = False

        # Bullets
        contenido_box = slide.shapes.add_textbox(
            PptxInches(0.6), PptxInches(1.7), ancho_texto, PptxInches(5.3)
        )
        tf = contenido_box.text_frame
        tf.word_wrap = True
        bullets = dia.get("bullets") or []
        if not bullets:
            bullets = [dia.get("contenido", "")] if dia.get("contenido") else []
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "•  " + str(b)
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_TEXTO
            p.space_after = Pt(14)

        # Notas del orador
        notas = dia.get("notas") or ""
        if notas:
            try:
                slide.notes_slide.notes_text_frame.text = notas
            except Exception:
                pass

        # Imagen o video a la derecha
        if hay_video:
            video_path = video_paths[video_idx]
            video_idx += 1
            try:
                with open(video_path, "rb") as vf:
                    video_bytes = BytesIO(vf.read())
                slide.shapes.add_movie(
                    video_bytes,
                    PptxInches(8.1), PptxInches(1.7),
                    PptxInches(4.6), PptxInches(3.5),
                    mime_type="video/mp4"
                )
            except Exception:
                print("Error insertando video en presentación:")
                traceback.print_exc()
        elif incluir_imagenes:
            img_bytes = imagenes_por_slide[idx]
            if img_bytes:
                try:
                    img_stream = BytesIO(img_bytes)
                    slide.shapes.add_picture(
                        img_stream, PptxInches(8.1), PptxInches(1.7), width=PptxInches(4.6)
                    )
                except Exception:
                    print("Error insertando imagen en presentación:")
                    traceback.print_exc()

    # ---- Videos sobrantes: se agregan como diapositivas extra ----
    while video_idx < len(video_paths):
        slide = prs.slides.add_slide(blank_layout)
        _agregar_fondo(slide, prs, COLOR_FONDO)

        titulo_box = slide.shapes.add_textbox(
            PptxInches(0.6), PptxInches(0.4), PptxInches(12.1), PptxInches(1)
        )
        p = titulo_box.text_frame.paragraphs[0]
        p.text = "Video"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITULO

        video_path = video_paths[video_idx]
        video_idx += 1
        try:
            with open(video_path, "rb") as vf:
                video_bytes = BytesIO(vf.read())
            slide.shapes.add_movie(
                video_bytes,
                PptxInches(2.0), PptxInches(1.5),
                PptxInches(9.3), PptxInches(5.6),
                mime_type="video/mp4"
            )
        except Exception:
            print("Error insertando video extra en presentación:")
            traceback.print_exc()

    nombre = f"presentacion_{uuid.uuid4().hex}.pptx"
    ruta = os.path.join(TEMP_DIR, nombre)
    prs.save(ruta)
    return ruta


def _procesar_presentacion_job(job_id, contenido_base, tema, titulo_pres, num_slides, incluir_imagenes, video_paths):
    """Corre en un hilo aparte: genera la estructura con IA, construye el .pptx
    y actualiza el estado del job. Al final borra los videos temporales subidos."""
    try:
        estructura = generar_estructura_presentacion(contenido_base, tema, num_slides)
        if not estructura:
            with PRESENTACIONES_LOCK:
                PRESENTACIONES_JOBS[job_id] = {
                    "status": "error",
                    "ruta": None,
                    "error": "No pude generar el contenido de la presentación. Probá de nuevo en unos segundos."
                }
            return

        if titulo_pres:
            estructura["titulo_presentacion"] = titulo_pres

        ruta_pptx = construir_pptx(estructura, incluir_imagenes=incluir_imagenes, video_paths=video_paths)

        with PRESENTACIONES_LOCK:
            PRESENTACIONES_JOBS[job_id] = {
                "status": "listo",
                "ruta": ruta_pptx,
                "error": None
            }

    except Exception as e:
        print("ERROR EN JOB DE PRESENTACION:", e)
        traceback.print_exc()
        with PRESENTACIONES_LOCK:
            PRESENTACIONES_JOBS[job_id] = {
                "status": "error",
                "ruta": None,
                "error": str(e)
            }
    finally:
        # limpiar videos temporales subidos
        for vp in (video_paths or []):
            try:
                if os.path.exists(vp):
                    os.remove(vp)
            except Exception:
                pass


@app.route("/generar_presentacion", methods=["POST"])
def generar_presentacion():
    """
    Inicia en segundo plano la generación de una presentación (.pptx) a partir
    de un tema descrito por el usuario y/o de un documento subido previamente
    (vía /upload_doc), con imágenes generadas por IA y/o videos cortos subidos
    por el usuario. Devuelve un job_id para consultar el progreso.
    """
    try:
        usuario = request.form.get("usuario_id", "anon")

        if not usuario_premium(usuario) and not es_superusuario(usuario):
            return jsonify({
                "ok": False,
                "error": "Esta función es exclusiva para usuarios Premium. Activá Foschi IA Premium para crear presentaciones."
            }), 403

        tema = (request.form.get("tema") or "").strip()
        titulo_pres = (request.form.get("titulo") or "").strip()
        doc_id = (request.form.get("doc_id") or "").strip()

        try:
            num_slides = int(request.form.get("num_slides", 8))
        except (TypeError, ValueError):
            num_slides = 8
        num_slides = max(3, min(num_slides, 20))

        incluir_imagenes = (request.form.get("incluir_imagenes", "true").lower() == "true")

        contenido_base = ""
        if doc_id:
            txt_path = os.path.join(TEMP_DIR, f"{doc_id}.txt")
            if os.path.exists(txt_path):
                try:
                    with open(txt_path, "r", encoding="utf-8") as f:
                        contenido_base = f.read()
                except Exception as e:
                    print("Error leyendo documento base para presentación:", e)

        if not contenido_base and not tema:
            return jsonify({
                "ok": False,
                "error": "Indicá un tema/descripción o subí un documento (📄 Analizar Documento) para generar la presentación."
            }), 400

        # Guardar videos subidos en disco YA (los FileStorage no sobreviven
        # al hilo en segundo plano)
        video_paths = []
        if "videos" in request.files:
            for v in request.files.getlist("videos"):
                if not v or not v.filename:
                    continue
                nombre_video = f"{uuid.uuid4().hex}_{secure_filename(v.filename)}"
                ruta_video = os.path.join(TEMP_DIR, nombre_video)
                try:
                    v.save(ruta_video)
                    video_paths.append(ruta_video)
                except Exception as e:
                    print("Error guardando video temporal:", e)

        job_id = uuid.uuid4().hex
        with PRESENTACIONES_LOCK:
            PRESENTACIONES_JOBS[job_id] = {"status": "procesando", "ruta": None, "error": None}

        hilo = threading.Thread(
            target=_procesar_presentacion_job,
            args=(job_id, contenido_base, tema, titulo_pres, num_slides, incluir_imagenes, video_paths),
            daemon=True
        )
        hilo.start()

        return jsonify({"ok": True, "job_id": job_id})

    except Exception as e:
        print("ERROR GENERAR PRESENTACION:", e)
        traceback.print_exc()
        return jsonify({"ok": False, "error": str(e)}), 500


@app.route("/estado_presentacion/<job_id>")
def estado_presentacion(job_id):
    with PRESENTACIONES_LOCK:
        job = PRESENTACIONES_JOBS.get(job_id)

    if not job:
        return jsonify({"ok": False, "error": "Job no encontrado"}), 404

    return jsonify({
        "ok": True,
        "status": job["status"],
        "error": job.get("error")
    })


@app.route("/descargar_presentacion/<job_id>")
def descargar_presentacion(job_id):
    with PRESENTACIONES_LOCK:
        job = PRESENTACIONES_JOBS.get(job_id)

    if not job or job["status"] != "listo" or not job.get("ruta"):
        return jsonify({"ok": False, "error": "La presentación todavía no está lista"}), 404

    ruta_pptx = job["ruta"]

    @after_this_request
    def _cleanup(response):
        try:
            if os.path.exists(ruta_pptx):
                os.remove(ruta_pptx)
        except Exception:
            pass
        with PRESENTACIONES_LOCK:
            PRESENTACIONES_JOBS.pop(job_id, None)
        return response

    return send_file(
        ruta_pptx,
        as_attachment=True,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        download_name="presentacion_foschi.pptx"
    )


# ---------------- RUN ----------------
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port)